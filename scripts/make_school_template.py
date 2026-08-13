#!/usr/bin/env python3
"""Generate examples/school-template/template.pptx — a reproducible
university-defense-style PPTX template with ZERO real school assets.

Style: navy (#1B2A4A) bands on white body, generic "××大学" wordmark text,
sample placeholder text on every fillable slot. Layout numbers are fixed so
the output is deterministic.

Run from repo root:  python scripts/make_school_template.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUTPUT = ROOT / "examples" / "school-template" / "template.pptx"

NAVY = RGBColor(0x1B, 0x2A, 0x4A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x4A, 0x55, 0x68)
PANEL = RGBColor(0xF2, 0xF4, 0xF8)
FONT = "Microsoft YaHei"


def _style_run(run, size_pt: float, color: RGBColor, *, bold: bool = False) -> None:
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
    # python-pptx only sets <a:latin>; CJK glyphs need <a:ea> as well.
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find(qn("a:ea"))
    if ea is None:
        ea = rpr.makeelement(qn("a:ea"), {})
        rpr.append(ea)
    ea.set("typeface", FONT)


def _textbox(slide, x, y, w, h, text, size_pt, color, *, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.word_wrap = True
    for index, line in enumerate(text.split("\n")):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = line
        _style_run(run, size_pt, color, bold=bold)
    return shape


def _rect(slide, x, y, w, h, fill: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _band_title(slide, text: str) -> None:
    """Navy title band with the section-title slot inside it."""
    band = _rect(slide, 0, 0, 13.333, 0.95, NAVY)
    frame = band.text_frame
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = text
    _style_run(run, 28, WHITE, bold=True)


def build_template(output: Path) -> Path:
    prs = Presentation()
    prs.slide_width = Emu(12192000)   # 13.333 in (16:9)
    prs.slide_height = Emu(6858000)   # 7.5 in
    blank = prs.slide_layouts[6]

    # Slide 1 — cover
    cover = prs.slides.add_slide(blank)
    _rect(cover, 0, 0, 13.333, 1.1, NAVY)
    _textbox(cover, 0.6, 0.3, 5.0, 0.55, "××大学", 20, WHITE, bold=True)
    _rect(cover, 0, 6.95, 13.333, 0.55, NAVY)
    _textbox(cover, 0.9, 2.7, 11.5, 1.3, "论文题目占位", 40, NAVY, bold=True, align=PP_ALIGN.CENTER)
    _textbox(cover, 0.9, 4.3, 11.5, 0.9, "答辩人：张三 | 指导教师：李四", 18, GRAY, align=PP_ALIGN.CENTER)

    # Slide 2 — table of contents
    toc = prs.slides.add_slide(blank)
    _rect(toc, 0.55, 0.72, 0.22, 0.5, NAVY)
    _textbox(toc, 0.95, 0.6, 4.0, 0.8, "目录", 32, NAVY, bold=True)
    _textbox(
        toc, 1.3, 1.9, 10.7, 5.0,
        "01  研究背景\n02  研究方法\n03  实验结果\n04  总结展望",
        20, GRAY,
    )

    # Slide 3 — content variant A (band title + body)
    content_a = prs.slides.add_slide(blank)
    _band_title(content_a, "点击输入章节标题")
    _textbox(content_a, 0.8, 1.6, 11.73, 5.3, "点击输入正文内容", 20, GRAY)

    # Slide 4 — content variant B (band title + light panel + body)
    content_b = prs.slides.add_slide(blank)
    _band_title(content_b, "点击输入章节标题")
    _rect(content_b, 0.6, 1.4, 12.13, 5.6, PANEL)
    _textbox(content_b, 0.9, 1.7, 11.53, 5.0, "点击输入正文内容", 20, GRAY)

    # Slide 5 — ending
    ending = prs.slides.add_slide(blank)
    _rect(ending, 0, 6.95, 13.333, 0.55, NAVY)
    _textbox(ending, 0.9, 2.7, 11.5, 1.2, "谢谢聆听", 40, NAVY, bold=True, align=PP_ALIGN.CENTER)
    _textbox(ending, 0.9, 4.2, 11.5, 0.7, "恳请各位老师批评指正", 20, GRAY, align=PP_ALIGN.CENTER)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    return output


if __name__ == "__main__":
    path = build_template(OUTPUT)
    print(path)
