from __future__ import annotations

import unittest
from xml.etree import ElementTree as ET

from slide_skill.filter_effects import (
    collect_filters,
    resolve_filter,
    apply_filter_to_shape,
    _normalize_color,
    _reorder_effect_lst,
    EFFECT_ORDER,
)
from slide_skill.svg_pipeline import (
    _soft_edge_filter_def,
    _glow_filter_def,
)


SVG_BLUR = """\
<svg xmlns="http://www.w3.org/2000/svg" width="960" height="540">
  <defs>
    <filter id="blur1">
      <feGaussianBlur stdDeviation="4"/>
    </filter>
  </defs>
  <rect x="10" y="10" width="200" height="100" filter="url(#blur1)"/>
</svg>
"""

SVG_SHADOW = """\
<svg xmlns="http://www.w3.org/2000/svg" width="960" height="540">
  <defs>
    <filter id="shadow1">
      <feDropShadow dx="2" dy="3" stdDeviation="4" flood-color="#000000" flood-opacity="0.25"/>
    </filter>
  </defs>
  <rect x="10" y="10" width="200" height="100" filter="url(#shadow1)"/>
</svg>
"""

SVG_CARD_SHADOW = """\
<svg xmlns="http://www.w3.org/2000/svg" width="960" height="540">
  <defs>
    <filter id="card-shadow-1" x="-5%" y="-5%" width="110%" height="110%">
      <feGaussianBlur in="SourceAlpha" stdDeviation="6" result="blur"/>
      <feOffset in="blur" dx="0" dy="3" result="offsetBlur"/>
      <feFlood flood-color="#000000" flood-opacity="0.16" result="shadowColor"/>
      <feComposite in="shadowColor" in2="offsetBlur" operator="in" result="shadow"/>
      <feMerge>
        <feMergeNode in="shadow"/>
        <feMergeNode in="SourceGraphic"/>
      </feMerge>
    </filter>
  </defs>
  <g filter="url(#card-shadow-1)">
    <rect x="10" y="10" width="200" height="100"/>
  </g>
</svg>
"""

SVG_NO_DEFS = """\
<svg xmlns="http://www.w3.org/2000/svg" width="960" height="540">
  <rect x="10" y="10" width="200" height="100"/>
</svg>
"""

SVG_SOFT_EDGE = """\
<svg xmlns="http://www.w3.org/2000/svg" width="960" height="540">
  <defs>
    <filter id="soft1" x="-20%" y="-20%" width="140%" height="150%">
      <feGaussianBlur in="SourceAlpha" stdDeviation="3"/>
    </filter>
  </defs>
  <rect x="10" y="10" width="200" height="100" filter="url(#soft1)"/>
</svg>
"""

SVG_GLOW = """\
<svg xmlns="http://www.w3.org/2000/svg" width="960" height="540">
  <defs>
    <filter id="glow1" x="-20%" y="-20%" width="140%" height="150%">
      <feGaussianBlur in="SourceAlpha" stdDeviation="4" result="blur"/>
      <feFlood flood-color="#FF6600" flood-opacity="0.5" result="glowColor"/>
      <feComposite in="glowColor" in2="blur" operator="in" result="glow"/>
      <feMerge>
        <feMergeNode in="glow"/>
        <feMergeNode in="SourceGraphic"/>
      </feMerge>
    </filter>
  </defs>
  <rect x="10" y="10" width="200" height="100" filter="url(#glow1)"/>
</svg>
"""

SVG_GLOW_WITH_OFFSET = """\
<svg xmlns="http://www.w3.org/2000/svg" width="960" height="540">
  <defs>
    <filter id="glow-offset" x="-20%" y="-20%" width="140%" height="150%">
      <feGaussianBlur in="SourceAlpha" stdDeviation="4" result="blur"/>
      <feOffset in="blur" dx="2" dy="3" result="offsetBlur"/>
      <feFlood flood-color="#FF6600" flood-opacity="0.5" result="glowColor"/>
      <feComposite in="glowColor" in2="offsetBlur" operator="in" result="glow"/>
      <feMerge>
        <feMergeNode in="glow"/>
        <feMergeNode in="SourceGraphic"/>
      </feMerge>
    </filter>
  </defs>
  <rect x="10" y="10" width="200" height="100" filter="url(#glow-offset)"/>
</svg>
"""

SVG_GLOW_AND_SHADOW = """\
<svg xmlns="http://www.w3.org/2000/svg" width="960" height="540">
  <defs>
    <filter id="glow-shadow" x="-20%" y="-20%" width="140%" height="150%">
      <feGaussianBlur in="SourceAlpha" stdDeviation="5" result="blur"/>
      <feOffset in="blur" dx="0" dy="4" result="offsetBlur"/>
      <feFlood flood-color="#3399FF" flood-opacity="0.6" result="glowColor"/>
      <feComposite in="glowColor" in2="blur" operator="in" result="glowResult"/>
      <feMerge>
        <feMergeNode in="glowResult"/>
        <feMergeNode in="SourceGraphic"/>
      </feMerge>
    </filter>
  </defs>
  <rect x="10" y="10" width="200" height="100" filter="url(#glow-shadow)"/>
</svg>
"""

SVG_SOURCE_GRAPHIC_BLUR = """\
<svg xmlns="http://www.w3.org/2000/svg" width="960" height="540">
  <defs>
    <filter id="sg-blur">
      <feGaussianBlur in="SourceGraphic" stdDeviation="4"/>
    </filter>
  </defs>
  <rect x="10" y="10" width="200" height="100" filter="url(#sg-blur)"/>
</svg>
"""

SVG_GLOW_DEFAULT_COLOR = """\
<svg xmlns="http://www.w3.org/2000/svg" width="960" height="540">
  <defs>
    <filter id="glow-default" x="-20%" y="-20%" width="140%" height="150%">
      <feGaussianBlur in="SourceAlpha" stdDeviation="4" result="blur"/>
      <feFlood flood-opacity="0.7" result="glowColor"/>
      <feComposite in="glowColor" in2="blur" operator="in"/>
    </filter>
  </defs>
  <rect x="10" y="10" width="200" height="100" filter="url(#glow-default)"/>
</svg>
"""


class CollectFiltersTest(unittest.TestCase):
    def test_blur_filter(self) -> None:
        root = ET.fromstring(SVG_BLUR)
        filters = collect_filters(root)
        self.assertIn("blur1", filters)
        f = filters["blur1"]
        self.assertIsNotNone(f["blur"])
        self.assertAlmostEqual(f["blur"]["stdDeviation"], 4.0)

    def test_shadow_filter(self) -> None:
        root = ET.fromstring(SVG_SHADOW)
        filters = collect_filters(root)
        self.assertIn("shadow1", filters)
        f = filters["shadow1"]
        self.assertIsNotNone(f["shadow"])
        self.assertAlmostEqual(f["shadow"]["dx"], 2.0)
        self.assertAlmostEqual(f["shadow"]["dy"], 3.0)
        self.assertAlmostEqual(f["shadow"]["stdDeviation"], 4.0)
        self.assertEqual(f["shadow"]["flood_color"], "000000")
        self.assertAlmostEqual(f["shadow"]["flood_opacity"], 0.25)

    def test_card_shadow_pattern(self) -> None:
        root = ET.fromstring(SVG_CARD_SHADOW)
        filters = collect_filters(root)
        self.assertIn("card-shadow-1", filters)
        f = filters["card-shadow-1"]
        self.assertIsNotNone(f["shadow"])
        self.assertAlmostEqual(f["shadow"]["dx"], 0.0)
        self.assertAlmostEqual(f["shadow"]["dy"], 3.0)
        self.assertAlmostEqual(f["shadow"]["stdDeviation"], 6.0)
        self.assertAlmostEqual(f["shadow"]["flood_opacity"], 0.16)

    def test_no_defs(self) -> None:
        root = ET.fromstring(SVG_NO_DEFS)
        filters = collect_filters(root)
        self.assertEqual(filters, {})

    def test_filter_no_id_skipped(self) -> None:
        svg = """\
<svg xmlns="http://www.w3.org/2000/svg" width="960" height="540">
  <defs>
    <filter>
      <feGaussianBlur stdDeviation="2"/>
    </filter>
  </defs>
</svg>"""
        root = ET.fromstring(svg)
        filters = collect_filters(root)
        self.assertEqual(filters, {})


class ResolveFilterTest(unittest.TestCase):
    def test_resolve_existing(self) -> None:
        root = ET.fromstring(SVG_BLUR)
        filters = collect_filters(root)
        result = resolve_filter("url(#blur1)", filters)
        self.assertIsNotNone(result)
        self.assertIsNotNone(result["blur"])

    def test_resolve_missing(self) -> None:
        root = ET.fromstring(SVG_BLUR)
        filters = collect_filters(root)
        result = resolve_filter("url(#nonexistent)", filters)
        self.assertIsNone(result)

    def test_resolve_not_url(self) -> None:
        result = resolve_filter("something", {})
        self.assertIsNone(result)


class NormalizeColorTest(unittest.TestCase):
    def test_hex6(self) -> None:
        self.assertEqual(_normalize_color("#FF0000"), "FF0000")

    def test_hex3(self) -> None:
        self.assertEqual(_normalize_color("#F00"), "FF0000")

    def test_no_hash(self) -> None:
        self.assertEqual(_normalize_color("00FF00"), "00FF00")

    def test_invalid(self) -> None:
        self.assertEqual(_normalize_color("xyz"), "000000")


class ApplyFilterTest(unittest.TestCase):
    def test_apply_blur_to_shape(self) -> None:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1))

        filter_info = {"blur": {"stdDeviation": 4.0}, "shadow": None}
        apply_filter_to_shape(shape, filter_info, 1.0, 1.0)

        from lxml import etree
        sp_pr = shape._element.spPr
        effect_elements = sp_pr.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst")
        self.assertGreaterEqual(len(effect_elements), 1)
        blur_elements = sp_pr.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}blur")
        self.assertGreaterEqual(len(blur_elements), 1)

    def test_apply_shadow_to_shape(self) -> None:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1))

        filter_info = {
            "blur": None,
            "shadow": {
                "dx": 2.0, "dy": 3.0, "stdDeviation": 4.0,
                "flood_color": "000000", "flood_opacity": 0.25,
            },
        }
        apply_filter_to_shape(shape, filter_info, 1.0, 1.0)

        from lxml import etree
        sp_pr = shape._element.spPr
        shdw_elements = sp_pr.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}outerShdw")
        self.assertGreaterEqual(len(shdw_elements), 1)

    def test_apply_no_effect_when_empty(self) -> None:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1))

        filter_info = {"blur": None, "shadow": None, "glow": None, "soft_edge": None}
        apply_filter_to_shape(shape, filter_info, 1.0, 1.0)

        sp_pr = shape._element.spPr
        effect_elements = sp_pr.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst")
        self.assertEqual(len(effect_elements), 0)

    def test_apply_soft_edge_to_shape(self) -> None:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1))

        filter_info = {
            "blur": None, "shadow": None,
            "soft_edge": {"stdDeviation": 3.0}, "glow": None,
        }
        apply_filter_to_shape(shape, filter_info, 1.0, 1.0)

        sp_pr = shape._element.spPr
        se_elements = sp_pr.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}softEdge")
        self.assertEqual(len(se_elements), 1)
        self.assertEqual(se_elements[0].get("rad"), str(int(3.0 * 25400)))

    def test_apply_glow_to_shape(self) -> None:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1))

        filter_info = {
            "blur": None, "shadow": None,
            "glow": {"stdDeviation": 4.0, "flood_color": "FF6600", "flood_opacity": 0.5},
            "soft_edge": None,
        }
        apply_filter_to_shape(shape, filter_info, 1.0, 1.0)

        sp_pr = shape._element.spPr
        glow_elements = sp_pr.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}glow")
        self.assertEqual(len(glow_elements), 1)
        self.assertEqual(glow_elements[0].get("rad"), str(int(4.0 * 25400)))
        srgb = glow_elements[0].find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr")
        self.assertIsNotNone(srgb)
        self.assertEqual(srgb.get("val"), "FF6600")
        alpha = srgb.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}alpha")
        self.assertIsNotNone(alpha)
        self.assertEqual(alpha.get("val"), str(int(0.5 * 100000)))

    def test_apply_glow_and_shadow_composition(self) -> None:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1))

        filter_info = {
            "blur": None,
            "shadow": {
                "dx": 0.0, "dy": 4.0, "stdDeviation": 5.0,
                "flood_color": "000000", "flood_opacity": 0.16,
            },
            "glow": {"stdDeviation": 5.0, "flood_color": "3399FF", "flood_opacity": 0.6},
            "soft_edge": None,
        }
        apply_filter_to_shape(shape, filter_info, 1.0, 1.0)

        sp_pr = shape._element.spPr
        glow_elements = sp_pr.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}glow")
        shdw_elements = sp_pr.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}outerShdw")
        self.assertGreaterEqual(len(glow_elements), 1)
        self.assertGreaterEqual(len(shdw_elements), 1)

    def test_effect_lst_xsd_ordering(self) -> None:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches
        from pptx.oxml.ns import qn

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1))

        filter_info = {
            "blur": None,
            "shadow": {
                "dx": 2.0, "dy": 3.0, "stdDeviation": 5.0,
                "flood_color": "000000", "flood_opacity": 0.25,
            },
            "glow": {"stdDeviation": 5.0, "flood_color": "3399FF", "flood_opacity": 0.6},
            "soft_edge": {"stdDeviation": 3.0},
        }
        apply_filter_to_shape(shape, filter_info, 1.0, 1.0)

        sp_pr = shape._element.spPr
        effectLst = sp_pr.find(qn("a:effectLst"))
        self.assertIsNotNone(effectLst)

        tags = [child.tag.split("}")[-1] for child in effectLst]
        expected_order = ["glow", "outerShdw", "softEdge"]
        actual_order = [t for t in tags if t in expected_order]
        self.assertEqual(actual_order, expected_order)

    def test_no_effect_when_all_none(self) -> None:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1))

        filter_info = {"blur": None, "shadow": None, "glow": None, "soft_edge": None}
        apply_filter_to_shape(shape, filter_info, 1.0, 1.0)

        sp_pr = shape._element.spPr
        effect_elements = sp_pr.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst")
        self.assertEqual(len(effect_elements), 0)


class SoftEdgeGlowDetectionTest(unittest.TestCase):
    def test_soft_edge_filter(self) -> None:
        root = ET.fromstring(SVG_SOFT_EDGE)
        filters = collect_filters(root)
        f = filters["soft1"]
        self.assertIsNotNone(f["soft_edge"])
        self.assertAlmostEqual(f["soft_edge"]["stdDeviation"], 3.0)
        self.assertIsNone(f["blur"])
        self.assertIsNone(f["glow"])
        self.assertIsNone(f["shadow"])

    def test_glow_filter(self) -> None:
        root = ET.fromstring(SVG_GLOW)
        filters = collect_filters(root)
        f = filters["glow1"]
        self.assertIsNotNone(f["glow"])
        self.assertAlmostEqual(f["glow"]["stdDeviation"], 4.0)
        self.assertEqual(f["glow"]["flood_color"], "FF6600")
        self.assertAlmostEqual(f["glow"]["flood_opacity"], 0.5)
        self.assertIsNone(f["shadow"])
        self.assertIsNone(f["soft_edge"])

    def test_glow_with_offset_classified_as_shadow(self) -> None:
        root = ET.fromstring(SVG_GLOW_WITH_OFFSET)
        filters = collect_filters(root)
        f = filters["glow-offset"]
        self.assertIsNotNone(f["shadow"])
        self.assertIsNone(f["glow"])

    def test_soft_edge_and_blur_mutually_exclusive(self) -> None:
        root = ET.fromstring(SVG_SOFT_EDGE)
        filters = collect_filters(root)
        f = filters["soft1"]
        self.assertIsNotNone(f["soft_edge"])
        self.assertIsNone(f["blur"])

    def test_glow_and_shadow_coexist(self) -> None:
        root = ET.fromstring(SVG_GLOW_AND_SHADOW)
        filters = collect_filters(root)
        f = filters["glow-shadow"]
        self.assertIsNotNone(f["glow"])
        self.assertIsNotNone(f["shadow"])

    def test_source_graphic_stays_blur(self) -> None:
        root = ET.fromstring(SVG_SOURCE_GRAPHIC_BLUR)
        filters = collect_filters(root)
        f = filters["sg-blur"]
        self.assertIsNotNone(f["blur"])
        self.assertIsNone(f["soft_edge"])
        self.assertIsNone(f["glow"])

    def test_unrecognized_pattern_falls_back_to_blur(self) -> None:
        svg = """\
<svg xmlns="http://www.w3.org/2000/svg" width="960" height="540">
  <defs>
    <filter id="unknown1">
      <feGaussianBlur stdDeviation="5"/>
    </filter>
  </defs>
</svg>"""
        root = ET.fromstring(svg)
        filters = collect_filters(root)
        f = filters["unknown1"]
        self.assertIsNotNone(f["blur"])
        self.assertIsNone(f["glow"])
        self.assertIsNone(f["soft_edge"])
        self.assertIsNone(f["shadow"])

    def test_glow_default_color(self) -> None:
        root = ET.fromstring(SVG_GLOW_DEFAULT_COLOR)
        filters = collect_filters(root)
        f = filters["glow-default"]
        self.assertIsNotNone(f["glow"])
        self.assertEqual(f["glow"]["flood_color"], "000000")
        self.assertAlmostEqual(f["glow"]["flood_opacity"], 0.7)


class ReorderEffectLstTest(unittest.TestCase):
    def test_ordering(self) -> None:
        from lxml import etree
        effectLst = etree.Element("effectLst")
        outerShdw = etree.SubElement(effectLst, "outerShdw")
        glow = etree.SubElement(effectLst, "glow")
        softEdge = etree.SubElement(effectLst, "softEdge")
        _reorder_effect_lst(effectLst)
        tags = [child.tag for child in effectLst]
        self.assertEqual(tags, ["glow", "outerShdw", "softEdge"])


class SvgPipelineFilterDefsTest(unittest.TestCase):
    def test_soft_edge_filter_def_output(self) -> None:
        result = _soft_edge_filter_def(1, std=5)
        self.assertIn('id="soft-edge-01"', result)
        self.assertIn('stdDeviation="5"', result)
        self.assertIn('in="SourceAlpha"', result)
        self.assertNotIn("feFlood", result)
        self.assertNotIn("feOffset", result)

    def test_glow_filter_def_output(self) -> None:
        result = _glow_filter_def(2, color="#FF6600", opacity=0.5, std=4)
        self.assertIn('id="glow-02"', result)
        self.assertIn('stdDeviation="4"', result)
        self.assertIn('flood-color="#FF6600"', result)
        self.assertIn('flood-opacity="0.5"', result)
        self.assertIn('operator="in"', result)
        self.assertIn("feMerge", result)

    def test_glow_filter_def_default_color(self) -> None:
        result = _glow_filter_def(3)
        self.assertIn('id="glow-03"', result)
        self.assertIn('flood-color="#000000"', result)


if __name__ == "__main__":
    unittest.main()
