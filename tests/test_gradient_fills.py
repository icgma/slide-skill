from __future__ import annotations

import tempfile
import unittest
from pathlib import Path
from xml.etree import ElementTree as ET

from slide_skill.gradient_fills import (
    apply_gradient_to_shape,
    collect_gradients,
    resolve_gradient_fill,
    _linear_angle,
    _percent,
)
from slide_skill.converters import _apply_fill_and_line


SVG_LINEAR = """\
<svg xmlns="http://www.w3.org/2000/svg" width="960" height="540">
  <defs>
    <linearGradient id="grad1" x1="0%" y1="0%" x2="100%" y2="0%">
      <stop offset="0%" stop-color="#FF0000"/>
      <stop offset="100%" stop-color="#0000FF"/>
    </linearGradient>
  </defs>
  <rect x="10" y="10" width="200" height="100" fill="url(#grad1)"/>
</svg>
"""

SVG_RADIAL = """\
<svg xmlns="http://www.w3.org/2000/svg" width="960" height="540">
  <defs>
    <radialGradient id="rgrad1" cx="50%" cy="50%" r="50%">
      <stop offset="0%" stop-color="#FFFFFF"/>
      <stop offset="100%" stop-color="#000000"/>
    </radialGradient>
  </defs>
  <circle cx="100" cy="100" r="50" fill="url(#rgrad1)"/>
</svg>
"""

SVG_NO_DEFS = """\
<svg xmlns="http://www.w3.org/2000/svg" width="960" height="540">
  <rect x="10" y="10" width="200" height="100" fill="#FF0000"/>
</svg>
"""

SVG_MULTIPLE = """\
<svg xmlns="http://www.w3.org/2000/svg" width="960" height="540">
  <defs>
    <linearGradient id="lg1" x1="0%" y1="0%" x2="100%" y2="0%">
      <stop offset="0%" stop-color="#AABBCC"/>
      <stop offset="50%" stop-color="#DDEEFF" stop-opacity="0.5"/>
      <stop offset="100%" stop-color="#112233"/>
    </linearGradient>
    <radialGradient id="rg1" cx="30%" cy="40%" r="60%">
      <stop offset="0%" stop-color="#FF0000"/>
      <stop offset="100%" stop-color="#00FF00"/>
    </radialGradient>
  </defs>
</svg>
"""


class CollectGradientsTest(unittest.TestCase):
    def test_linear_gradient(self) -> None:
        root = ET.fromstring(SVG_LINEAR)
        grads = collect_gradients(root)
        self.assertIn("grad1", grads)
        g = grads["grad1"]
        self.assertEqual(g["type"], "linear")
        self.assertEqual(len(g["stops"]), 2)
        self.assertAlmostEqual(g["stops"][0]["offset"], 0.0)
        self.assertAlmostEqual(g["stops"][1]["offset"], 1.0)

    def test_radial_gradient(self) -> None:
        root = ET.fromstring(SVG_RADIAL)
        grads = collect_gradients(root)
        self.assertIn("rgrad1", grads)
        g = grads["rgrad1"]
        self.assertEqual(g["type"], "radial")
        self.assertAlmostEqual(g["cx"], 0.5)
        self.assertAlmostEqual(g["cy"], 0.5)

    def test_no_defs(self) -> None:
        root = ET.fromstring(SVG_NO_DEFS)
        grads = collect_gradients(root)
        self.assertEqual(grads, {})

    def test_multiple_gradients(self) -> None:
        root = ET.fromstring(SVG_MULTIPLE)
        grads = collect_gradients(root)
        self.assertIn("lg1", grads)
        self.assertIn("rg1", grads)
        self.assertEqual(grads["lg1"]["type"], "linear")
        self.assertEqual(grads["rg1"]["type"], "radial")
        self.assertEqual(len(grads["lg1"]["stops"]), 3)
        self.assertAlmostEqual(grads["lg1"]["stops"][1]["opacity"], 0.5)

    def test_gradient_with_no_id_skipped(self) -> None:
        svg = """\
<svg xmlns="http://www.w3.org/2000/svg" width="960" height="540">
  <defs>
    <linearGradient x1="0%" y1="0%" x2="100%" y2="0%">
      <stop offset="0%" stop-color="#FF0000"/>
    </linearGradient>
  </defs>
</svg>"""
        root = ET.fromstring(svg)
        grads = collect_gradients(root)
        self.assertEqual(grads, {})


class ResolveGradientFillTest(unittest.TestCase):
    def test_resolve_linear(self) -> None:
        root = ET.fromstring(SVG_LINEAR)
        grads = collect_gradients(root)
        result = resolve_gradient_fill("url(#grad1)", grads)
        self.assertIsNotNone(result)
        self.assertEqual(result["type"], "linear")

    def test_resolve_missing(self) -> None:
        root = ET.fromstring(SVG_LINEAR)
        grads = collect_gradients(root)
        result = resolve_gradient_fill("url(#nonexistent)", grads)
        self.assertIsNone(result)

    def test_resolve_not_url(self) -> None:
        result = resolve_gradient_fill("#FF0000", {})
        self.assertIsNone(result)


class LinearAngleTest(unittest.TestCase):
    def test_horizontal_left_to_right(self) -> None:
        angle = _linear_angle(0.0, 0.0, 1.0, 0.0)
        self.assertAlmostEqual(angle, 90.0, places=1)

    def test_vertical_top_to_bottom(self) -> None:
        angle = _linear_angle(0.0, 0.0, 0.0, 1.0)
        self.assertAlmostEqual(angle, 0.0, places=1)

    def test_diagonal(self) -> None:
        angle = _linear_angle(0.0, 0.0, 1.0, 1.0)
        self.assertAlmostEqual(angle, 45.0, places=1)

    def test_degenerate(self) -> None:
        angle = _linear_angle(0.5, 0.5, 0.5, 0.5)
        self.assertAlmostEqual(angle, 0.0, places=1)


class PercentTest(unittest.TestCase):
    def test_percent_value(self) -> None:
        self.assertAlmostEqual(_percent("50%"), 0.5)

    def test_fraction_value(self) -> None:
        self.assertAlmostEqual(_percent("0.75"), 0.75)

    def test_zero(self) -> None:
        self.assertAlmostEqual(_percent("0%"), 0.0)

    def test_invalid(self) -> None:
        self.assertAlmostEqual(_percent("abc"), 0.0)


class ApplyGradientToShapeTest(unittest.TestCase):
    def test_apply_linear_to_shape(self) -> None:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1))

        gradient = {
            "type": "linear",
            "angle": 90.0,
            "stops": [
                {"offset": 0.0, "color": "#FF0000", "opacity": 1.0},
                {"offset": 1.0, "color": "#0000FF", "opacity": 1.0},
            ],
        }
        apply_gradient_to_shape(shape, gradient, RGBColor)

        from lxml import etree
        sp_pr = shape._element.spPr
        gfill_elements = sp_pr.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}gradFill")
        self.assertGreaterEqual(len(gfill_elements), 1)

    def test_apply_radial_to_shape(self) -> None:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(1), Inches(2), Inches(2))

        gradient = {
            "type": "radial",
            "cx": 0.5, "cy": 0.5, "r": 0.5,
            "fx": 0.5, "fy": 0.5,
            "stops": [
                {"offset": 0.0, "color": "#FFFFFF", "opacity": 1.0},
                {"offset": 1.0, "color": "#000000", "opacity": 1.0},
            ],
        }
        apply_gradient_to_shape(shape, gradient, RGBColor)

        from lxml import etree
        sp_pr = shape._element.spPr
        gfill_elements = sp_pr.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}gradFill")
        self.assertGreaterEqual(len(gfill_elements), 1)
        path_elements = sp_pr.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}path")
        self.assertGreaterEqual(len(path_elements), 1)


class ConverterGradientIntegrationTest(unittest.TestCase):
    def test_apply_fill_and_line_with_gradient(self) -> None:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1))

        svg_with_grad = (
            '<svg xmlns="http://www.w3.org/2000/svg" width="960" height="540">'
            '<defs><linearGradient id="mygrad" x1="0%" y1="0%" x2="0%" y2="100%">'
            '<stop offset="0%" stop-color="#FF0000"/>'
            '<stop offset="100%" stop-color="#00FF00"/>'
            '</linearGradient></defs>'
            '<rect fill="url(#mygrad)"/></svg>'
        )
        root = ET.fromstring(svg_with_grad)
        elem = root.find(".//{http://www.w3.org/2000/svg}rect")
        _apply_fill_and_line(shape, elem, RGBColor, root=root)

        sp_pr = shape._element.spPr
        gfill_elements = sp_pr.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}gradFill")
        self.assertGreaterEqual(len(gfill_elements), 1)

    def test_apply_fill_and_line_solid_when_no_gradient(self) -> None:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1))

        elem = ET.Element("rect", attrib={"fill": "#FF0000"})
        _apply_fill_and_line(shape, elem, RGBColor)
        self.assertEqual(str(shape.fill.fore_color.rgb), "FF0000")

    def test_apply_fill_and_line_url_without_gradients_dict(self) -> None:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1))

        elem = ET.Element("rect", attrib={"fill": "url(#mygrad)"})
        _apply_fill_and_line(shape, elem, RGBColor)
        sp_pr = shape._element.spPr
        gfill_elements = sp_pr.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}gradFill")
        self.assertEqual(len(gfill_elements), 0)


class ExportWithGradientsTest(unittest.TestCase):
    def test_export_svg_with_linear_gradient(self) -> None:
        from slide_skill.exporter import export_project, validate_pptx
        from slide_skill.project import init_project
        from slide_skill.svg_pipeline import finalize_svg, generate_svg

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "source.md"
            source.write_text("# Gradient Test\n\n## Slide 1\n\nGradient fill test.\n", encoding="utf-8")
            project = init_project("GradTest", base_dir=root / "projects")
            from slide_skill.svg_pipeline import create_spec
            create_spec(project, source)
            svg_paths = generate_svg(project, source)
            self.assertGreaterEqual(len(svg_paths), 1)

            svg_file = svg_paths[0]
            svg_content = svg_file.read_text(encoding="utf-8")
            if "<defs>" not in svg_content:
                svg_with_grad = svg_content.replace(
                    "<svg ",
                    '<svg xmlns="http://www.w3.org/2000/svg" '
                )
                insert_pos = svg_with_grad.find(">", svg_with_grad.find("<svg")) + 1
                grad_defs = (
                    '<defs><linearGradient id="testGrad" x1="0%" y1="0%" x2="100%" y2="0%">'
                    '<stop offset="0%" stop-color="#FF0000"/>'
                    '<stop offset="100%" stop-color="#0000FF"/>'
                    '</linearGradient></defs>'
                )
                svg_with_grad = svg_with_grad[:insert_pos] + grad_defs + svg_with_grad[insert_pos:]
                svg_file.write_text(svg_with_grad, encoding="utf-8")

            from slide_skill.svg_pipeline import check_project_svg
            ok, _ = check_project_svg(project, stage="output")

            final_paths = finalize_svg(project)
            self.assertGreaterEqual(len(final_paths), 1)
            deck = export_project(project)
            valid, errors = validate_pptx(deck)
            self.assertTrue(valid, errors)


if __name__ == "__main__":
    unittest.main()
