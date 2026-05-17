from __future__ import annotations

import unittest
from pathlib import Path
from unittest.mock import patch, MagicMock

from slide_skill.pdf_handout import (
    _parse_notes_markdown,
    _page_aspect,
    Layout,
)


class ParseNotesMarkdownTest(unittest.TestCase):
    def test_plain_text(self) -> None:
        segments = _parse_notes_markdown("Hello world")
        self.assertEqual(len(segments), 1)
        self.assertEqual(segments[0], ("Hello world", False, False, False))

    def test_bold(self) -> None:
        segments = _parse_notes_markdown("**bold text**")
        self.assertTrue(any(s[1] for s in segments))

    def test_italic(self) -> None:
        segments = _parse_notes_markdown("*italic text*")
        self.assertTrue(any(s[2] for s in segments))

    def test_bullet(self) -> None:
        segments = _parse_notes_markdown("- bullet item")
        self.assertTrue(any(s[3] for s in segments))

    def test_mixed_formatting(self) -> None:
        text = "Normal **bold** and *italic* text"
        segments = _parse_notes_markdown(text)
        has_normal = any(not s[1] and not s[2] and s[0] for s in segments)
        has_bold = any(s[1] for s in segments)
        has_italic = any(s[2] for s in segments)
        self.assertTrue(has_normal)
        self.assertTrue(has_bold)
        self.assertTrue(has_italic)

    def test_empty_string(self) -> None:
        segments = _parse_notes_markdown("")
        self.assertEqual(len(segments), 1)
        self.assertEqual(segments[0], ("", False, False, False))

    def test_multiple_lines(self) -> None:
        text = "Line one\nLine two\n- Bullet"
        segments = _parse_notes_markdown(text)
        self.assertEqual(len(segments), 3)

    def test_bullet_with_bold(self) -> None:
        segments = _parse_notes_markdown("- **bold bullet**")
        bullet_bold = [s for s in segments if s[1] and s[3]]
        self.assertTrue(bullet_bold)


class PageAspectTest(unittest.TestCase):
    def test_ppt169_landscape(self) -> None:
        w, h = _page_aspect("ppt169")
        self.assertGreater(w, h)

    def test_xhs_portrait(self) -> None:
        w, h = _page_aspect("xhs")
        self.assertLess(w, h)

    def test_square_aspect(self) -> None:
        w, h = _page_aspect("square")
        self.assertAlmostEqual(w, h, delta=1)

    def test_unknown_format_fallback(self) -> None:
        w, h = _page_aspect("nonexistent")
        self.assertGreater(w, 0)
        self.assertGreater(h, 0)

    def test_story_portrait(self) -> None:
        w, h = _page_aspect("story")
        self.assertLess(w, h)


class ExportHandoutIntegrationTest(unittest.TestCase):
    """Integration tests that mock PPTX/PDF generation to avoid soffice dependency."""

    @patch("slide_skill.pdf_handout._export_thumbnails")
    @patch("slide_skill.project.load_project")
    def test_1up_generates_pdf(
        self, mock_load, mock_thumbs
    ) -> None:
        import tempfile
        from PIL import Image

        mock_load.return_value = {"name": "test"}

        # Create a minimal PPTX to avoid "no exports" error
        with tempfile.TemporaryDirectory() as tmpdir:
            project = Path(tmpdir)
            exports = project / "exports"
            exports.mkdir()

            # Create minimal PPTX
            from pptx import Presentation
            prs = Presentation()
            prs.slides.add_slide(prs.slide_layouts[6])
            pptx_path = exports / "test.pptx"
            prs.save(str(pptx_path))

            # Create placeholder thumb
            thumb_dir = project / ".thumbs"
            thumb_dir.mkdir()
            thumb = thumb_dir / "slide_000.png"
            img = Image.new("RGB", (320, 180), "#EEEEEE")
            img.save(str(thumb))

            mock_thumbs.return_value = [thumb]
            mock_load.return_value = {"name": "test"}

            output = exports / "handout.pdf"
            from slide_skill.pdf_handout import export_handout
            result = export_handout(project, output, layout="1-up", canvas="ppt169",
                                   notes=["Test note with **bold** text"])

            self.assertTrue(result.exists())
            self.assertGreater(result.stat().st_size, 0)

    @patch("slide_skill.pdf_handout._export_thumbnails")
    @patch("slide_skill.project.load_project")
    def test_2up_generates_pdf(
        self, mock_load, mock_thumbs
    ) -> None:
        import tempfile
        from PIL import Image

        mock_load.return_value = {"name": "test2"}

        with tempfile.TemporaryDirectory() as tmpdir:
            project = Path(tmpdir)
            exports = project / "exports"
            exports.mkdir()

            from pptx import Presentation
            prs = Presentation()
            prs.slides.add_slide(prs.slide_layouts[6])
            prs.slides.add_slide(prs.slide_layouts[6])
            pptx_path = exports / "test2.pptx"
            prs.save(str(pptx_path))

            thumb_dir = project / ".thumbs"
            thumb_dir.mkdir()
            thumbs = []
            for i in range(2):
                t = thumb_dir / f"slide_{i:03d}.png"
                img = Image.new("RGB", (320, 180), "#EEEEEE")
                img.save(str(t))
                thumbs.append(t)

            mock_thumbs.return_value = thumbs
            mock_load.return_value = {"name": "test2"}

            output = exports / "handout2.pdf"
            from slide_skill.pdf_handout import export_handout
            result = export_handout(project, output, layout="2-up", canvas="ppt169",
                                   notes=["Note one", "Note two"])

            self.assertTrue(result.exists())

    @patch("slide_skill.pdf_handout._export_thumbnails")
    @patch("slide_skill.project.load_project")
    def test_cjk_notes_render(
        self, mock_load, mock_thumbs
    ) -> None:
        import tempfile
        from PIL import Image
        from slide_skill.pdf_handout import _add_cjk_font
        from fpdf import FPDF

        # Skip if no CJK font available
        test_pdf = FPDF()
        if not _add_cjk_font(test_pdf):
            self.skipTest("No CJK font available on this system")

        mock_load.return_value = {"name": "test_cjk"}

        with tempfile.TemporaryDirectory() as tmpdir:
            project = Path(tmpdir)
            exports = project / "exports"
            exports.mkdir()

            from pptx import Presentation
            prs = Presentation()
            prs.slides.add_slide(prs.slide_layouts[6])
            pptx_path = exports / "test_cjk.pptx"
            prs.save(str(pptx_path))

            thumb_dir = project / ".thumbs"
            thumb_dir.mkdir()
            thumb = thumb_dir / "slide_000.png"
            img = Image.new("RGB", (320, 180), "#EEEEEE")
            img.save(str(thumb))

            mock_thumbs.return_value = [thumb]
            mock_load.return_value = {"name": "test_cjk"}

            output = exports / "cjk_handout.pdf"
            from slide_skill.pdf_handout import export_handout
            result = export_handout(project, output, layout="1-up", canvas="ppt169",
                                   notes=["中文备注\n- 项目一\n*斜体*"])

            self.assertTrue(result.exists())
            self.assertGreater(result.stat().st_size, 0)

    @patch("slide_skill.pdf_handout._export_thumbnails")
    @patch("slide_skill.project.load_project")
    def test_3up_with_canvas_aspect(
        self, mock_load, mock_thumbs
    ) -> None:
        import tempfile
        from PIL import Image

        mock_load.return_value = {"name": "test_3up"}

        with tempfile.TemporaryDirectory() as tmpdir:
            project = Path(tmpdir)
            exports = project / "exports"
            exports.mkdir()

            from pptx import Presentation
            prs = Presentation()
            for _ in range(3):
                prs.slides.add_slide(prs.slide_layouts[6])
            pptx_path = exports / "test_3up.pptx"
            prs.save(str(pptx_path))

            thumb_dir = project / ".thumbs"
            thumb_dir.mkdir()
            thumbs = []
            for i in range(3):
                t = thumb_dir / f"slide_{i:03d}.png"
                img = Image.new("RGB", (320, 180), "#EEEEEE")
                img.save(str(t))
                thumbs.append(t)

            mock_thumbs.return_value = thumbs
            mock_load.return_value = {"name": "test_3up"}

            output = exports / "handout3.pdf"
            from slide_skill.pdf_handout import export_handout
            result = export_handout(
                project, output, layout="3-up", canvas="xhs",
                notes=["Note 1", "Note 2", "Note 3"])

            self.assertTrue(result.exists())


if __name__ == "__main__":
    unittest.main()
