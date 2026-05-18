from __future__ import annotations

import unittest
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from slide_skill.converters import (
    convert_bilingual_group,
    approx_w_in,
    _parse_bilingual_font_sizes,
)
from slide_skill.i18n import CJK_CHAR_WIDTH, LATIN_CHAR_WIDTH, detect_language


SVG_BILINGUAL_STACKED = """\
<svg xmlns="http://www.w3.org/2000/svg" width="960" height="540">
  <g data-bilingual="true">
    <text x="80" y="60" font-size="24" fill="#333333">中文标题</text>
    <text x="80" y="90" font-size="18" fill="#666666">English Title</text>
  </g>
</svg>
"""

SVG_BILINGUAL_SIDE_BY_SIDE = """\
<svg xmlns="http://www.w3.org/2000/svg" width="960" height="540">
  <g data-bilingual="true" data-layout="side-by-side">
    <text x="80" y="60" font-size="24" fill="#333333">中文标题</text>
    <text x="80" y="90" font-size="18" fill="#666666">English Title</text>
  </g>
</svg>
"""

SVG_BILINGUAL_CUSTOM_SIZE = """\
<svg xmlns="http://www.w3.org/2000/svg" width="960" height="540">
  <g data-bilingual="true" data-lang-size="zh:20;en:14">
    <text x="80" y="60" font-size="24" fill="#333333">中文标题</text>
    <text x="80" y="90" font-size="18" fill="#666666">English Title</text>
  </g>
</svg>
"""

SVG_BILINGUAL_EN_ONLY_SIZE = """\
<svg xmlns="http://www.w3.org/2000/svg" width="960" height="540">
  <g data-bilingual="true" data-lang-size="16">
    <text x="80" y="60" font-size="24" fill="#333333">中文标题</text>
    <text x="80" y="90" font-size="18" fill="#666666">English Title</text>
  </g>
</svg>
"""

SVG_BILINGUAL_SINGLE_TEXT = """\
<svg xmlns="http://www.w3.org/2000/svg" width="960" height="540">
  <g data-bilingual="true">
    <text x="80" y="60" font-size="24">Only one text</text>
  </g>
</svg>
"""

SVG_BILINGUAL_SAME_LANG = """\
<svg xmlns="http://www.w3.org/2000/svg" width="960" height="540">
  <g data-bilingual="true">
    <text x="80" y="60" font-size="24">English One</text>
    <text x="80" y="90" font-size="18">English Two</text>
  </g>
</svg>
"""


def _make_slide():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    scale_x = 13.33 / 960
    scale_y = 7.5 / 540
    meta = {
        "canvas": {
            "pptx_width_in": "13.33",
            "pptx_height_in": "7.5",
            "width": "960",
            "height": "540",
        }
    }
    return slide, scale_x, scale_y, meta


class BilingualDetectionTest(unittest.TestCase):
    def test_stacked_produces_two_text_frames(self) -> None:
        slide, sx, sy, meta = _make_slide()
        root = ET.fromstring(SVG_BILINGUAL_STACKED)
        group = root.find(".//{http://www.w3.org/2000/svg}g")
        shapes = convert_bilingual_group(slide, group, sx, sy, meta, RGBColor)
        self.assertEqual(len(shapes), 2)
        zh_text = shapes[0].text_frame.text
        en_text = shapes[1].text_frame.text
        self.assertIn("中文", zh_text)
        self.assertIn("English", en_text)

    def test_side_by_side_produces_two_text_frames(self) -> None:
        slide, sx, sy, meta = _make_slide()
        root = ET.fromstring(SVG_BILINGUAL_SIDE_BY_SIDE)
        group = root.find(".//{http://www.w3.org/2000/svg}g")
        shapes = convert_bilingual_group(slide, group, sx, sy, meta, RGBColor)
        self.assertEqual(len(shapes), 2)
        # English frame should be to the right of Chinese frame
        zh_left = shapes[0].left
        en_left = shapes[1].left
        self.assertGreater(en_left, zh_left)

    def test_english_smaller_font(self) -> None:
        slide, sx, sy, meta = _make_slide()
        root = ET.fromstring(SVG_BILINGUAL_STACKED)
        group = root.find(".//{http://www.w3.org/2000/svg}g")
        shapes = convert_bilingual_group(slide, group, sx, sy, meta, RGBColor)
        zh_size = shapes[0].text_frame.paragraphs[0].runs[0].font.size
        en_size = shapes[1].text_frame.paragraphs[0].runs[0].font.size
        self.assertGreater(zh_size, en_size)

    def test_custom_lang_size_zh_en(self) -> None:
        slide, sx, sy, meta = _make_slide()
        root = ET.fromstring(SVG_BILINGUAL_CUSTOM_SIZE)
        group = root.find(".//{http://www.w3.org/2000/svg}g")
        shapes = convert_bilingual_group(slide, group, sx, sy, meta, RGBColor)
        # zh:20pt, en:14pt
        from pptx.util import Pt
        zh_size = shapes[0].text_frame.paragraphs[0].runs[0].font.size
        en_size = shapes[1].text_frame.paragraphs[0].runs[0].font.size
        self.assertEqual(zh_size, Pt(20))
        self.assertEqual(en_size, Pt(14))

    def test_en_only_size_attribute(self) -> None:
        slide, sx, sy, meta = _make_slide()
        root = ET.fromstring(SVG_BILINGUAL_EN_ONLY_SIZE)
        group = root.find(".//{http://www.w3.org/2000/svg}g")
        shapes = convert_bilingual_group(slide, group, sx, sy, meta, RGBColor)
        from pptx.util import Pt
        en_size = shapes[1].text_frame.paragraphs[0].runs[0].font.size
        self.assertEqual(en_size, Pt(16))

    def test_single_text_returns_empty(self) -> None:
        slide, sx, sy, meta = _make_slide()
        root = ET.fromstring(SVG_BILINGUAL_SINGLE_TEXT)
        group = root.find(".//{http://www.w3.org/2000/svg}g")
        shapes = convert_bilingual_group(slide, group, sx, sy, meta, RGBColor)
        self.assertEqual(len(shapes), 0)

    def test_same_language_returns_empty(self) -> None:
        slide, sx, sy, meta = _make_slide()
        root = ET.fromstring(SVG_BILINGUAL_SAME_LANG)
        group = root.find(".//{http://www.w3.org/2000/svg}g")
        shapes = convert_bilingual_group(slide, group, sx, sy, meta, RGBColor)
        self.assertEqual(len(shapes), 0)

    def test_stacked_english_below_chinese(self) -> None:
        slide, sx, sy, meta = _make_slide()
        root = ET.fromstring(SVG_BILINGUAL_STACKED)
        group = root.find(".//{http://www.w3.org/2000/svg}g")
        shapes = convert_bilingual_group(slide, group, sx, sy, meta, RGBColor)
        zh_top = shapes[0].top
        en_top = shapes[1].top
        self.assertGreater(en_top, zh_top)


class WidthEstimationTest(unittest.TestCase):
    def test_cjk_wider_than_latin(self) -> None:
        cjk_w = approx_w_in("中文标题", 24.0)
        latin_w = approx_w_in("abcd", 24.0)
        self.assertGreater(cjk_w, latin_w)

    def test_uses_i18n_constants(self) -> None:
        self.assertEqual(CJK_CHAR_WIDTH, 1.0)
        self.assertEqual(LATIN_CHAR_WIDTH, 0.55)

    def test_empty_string_min_width(self) -> None:
        w = approx_w_in("", 24.0)
        self.assertGreater(w, 0)


class ParseBilingualFontSizesTest(unittest.TestCase):
    def test_default_ratio(self) -> None:
        zh, en = _parse_bilingual_font_sizes("", 24.0)
        self.assertEqual(zh, 24.0)
        self.assertAlmostEqual(en, 18.0)

    def test_custom_zh_en(self) -> None:
        zh, en = _parse_bilingual_font_sizes("zh:20;en:14", 24.0)
        self.assertEqual(zh, 20.0)
        self.assertEqual(en, 14.0)

    def test_single_value(self) -> None:
        zh, en = _parse_bilingual_font_sizes("16", 24.0)
        self.assertEqual(zh, 24.0)
        self.assertEqual(en, 16.0)

    def test_pt_suffix(self) -> None:
        zh, en = _parse_bilingual_font_sizes("16pt", 24.0)
        self.assertEqual(en, 16.0)

    def test_invalid_returns_default(self) -> None:
        zh, en = _parse_bilingual_font_sizes("abc", 24.0)
        self.assertEqual(zh, 24.0)
        self.assertAlmostEqual(en, 18.0)


class LanguageDetectionTest(unittest.TestCase):
    def test_detects_chinese(self) -> None:
        self.assertEqual(detect_language("中文标题"), "zh")

    def test_detects_english(self) -> None:
        self.assertEqual(detect_language("English text"), "en")

    def test_detects_arabic(self) -> None:
        self.assertEqual(detect_language("مرحبا"), "ar")

    def test_detects_hebrew(self) -> None:
        self.assertEqual(detect_language("שלום"), "he")

    def test_detects_japanese(self) -> None:
        self.assertEqual(detect_language("こんにちは"), "ja")

    def test_detects_korean(self) -> None:
        self.assertEqual(detect_language("안녕하세요"), "ko")


if __name__ == "__main__":
    unittest.main()
