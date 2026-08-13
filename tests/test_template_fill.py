from __future__ import annotations

import json
import tempfile
import unittest
import zipfile
from pathlib import Path

from slide_skill import cli
from slide_skill.intake import extract_pptx_slide_text
from slide_skill.template_fill import (
    analyze_template,
    fill_template,
    parse_content_markdown,
)


def _build_template(path: Path, *, planted_placeholder: bool = False) -> None:
    """Synthetic 5-slide university-style template (python-pptx for fixture only)."""
    from pptx import Presentation
    from pptx.util import Emu, Inches, Pt

    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    blank = prs.slide_layouts[6]

    def box(slide, x, y, w, h, text, size_pt):
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        frame = shape.text_frame
        frame.word_wrap = True
        for index, line in enumerate(text.split("\n")):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            run = paragraph.add_run()
            run.text = line
            run.font.size = Pt(size_pt)
        return shape

    cover = prs.slides.add_slide(blank)
    box(cover, 1, 2.5, 11, 1.2, "论文题目样板", 40)
    box(cover, 1, 4.2, 11, 0.8, "答辩人：张三 | 指导教师：李四", 18)
    toc = prs.slides.add_slide(blank)
    box(toc, 1, 0.5, 4, 0.8, "目录", 32)
    box(toc, 1, 1.8, 11, 4.5, "一、样例章节\n二、样例章节", 20)
    for _ in range(2):
        page = prs.slides.add_slide(blank)
        box(page, 0.6, 0.3, 11.5, 0.8, "点击输入章节标题", 28)
        box(page, 0.8, 1.5, 11.5, 1.6, "点击输入正文内容", 20)
    if planted_placeholder:
        box(prs.slides[3], 0.8, 5.8, 6, 0.6, "点击输入标题", 14)
    ending = prs.slides.add_slide(blank)
    box(ending, 1, 2.8, 11, 1.2, "谢谢聆听", 40)
    prs.save(path)


_FOUR_SECTION_MD = """# 测试论文标题

> 答辩人:李雷 | 指导教师:韩梅梅 | 2026 年 6 月

## 研究背景

- 背景要点一
- 背景要点二

## 研究方法

- 方法要点一
- 方法要点二

## 实验结果

- 结果要点一
- 结果要点二

## 结论展望

- 结论要点一
- 结论要点二
"""


class ParseContentMarkdownTest(unittest.TestCase):
    def test_extracts_title_meta_and_sections(self) -> None:
        content = parse_content_markdown(_FOUR_SECTION_MD)
        self.assertEqual("测试论文标题", content.title)
        self.assertEqual(["答辩人:李雷 | 指导教师:韩梅梅 | 2026 年 6 月"], content.meta_lines)
        self.assertEqual(4, len(content.sections))
        self.assertEqual("研究背景", content.sections[0].heading)
        self.assertEqual(["背景要点一", "背景要点二"], content.sections[0].body_lines)
        self.assertEqual([], content.toc_lines)

    def test_toc_like_section_feeds_toc_not_content(self) -> None:
        md = "# 标题\n\n## 汇报提纲\n\n- 第一部分\n- 第二部分\n\n## 正文章节\n\n- 要点\n"
        content = parse_content_markdown(md)
        self.assertEqual(["第一部分", "第二部分"], content.toc_lines)
        self.assertEqual(["正文章节"], [section.heading for section in content.sections])

    def test_table_rows_become_body_lines(self) -> None:
        md = "# 标题\n\n## 参数设置\n\n| 参数 | 取值 |\n|---|---|\n| 窗口 | 14 天 |\n"
        content = parse_content_markdown(md)
        self.assertEqual(["参数 | 取值", "窗口 | 14 天"], content.sections[0].body_lines)


class TemplateFillTest(unittest.TestCase):
    def test_fill_expands_pages_and_preserves_format(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            template = root / "template.pptx"
            _build_template(template)
            content = root / "thesis.md"
            content.write_text(_FOUR_SECTION_MD, encoding="utf-8")
            output = root / "filled.pptx"

            result = fill_template(template, content, output)

            slides = extract_pptx_slide_text(output)
            self.assertEqual(7, len(slides))  # cover + toc + 4 content + ending
            self.assertIn("测试论文标题", slides[0][1])
            self.assertNotIn("论文题目样板", slides[0][1])
            self.assertIn("李雷", slides[0][1])
            self.assertNotIn("张三", slides[0][1])
            self.assertIn("1. 研究背景", slides[1][1])
            self.assertIn("4. 结论展望", slides[1][1])
            for offset, heading in enumerate(["研究背景", "研究方法", "实验结果", "结论展望"]):
                self.assertIn(heading, slides[2 + offset][1])
            self.assertIn("背景要点二", slides[2][1])
            self.assertIn("谢谢聆听", slides[6][1])
            # Template formatting survives: the 40pt cover-title run keeps its sz attr.
            with zipfile.ZipFile(output) as zf:
                cover_xml = zf.read("ppt/slides/slide1.xml").decode("utf-8")
            self.assertIn('sz="4000"', cover_xml)
            self.assertEqual("clean", result.verdict)
            self.assertEqual([], result.stale)
            self.assertEqual([], result.overflow)

    def test_surplus_content_slides_deleted(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            template = root / "template.pptx"
            _build_template(template)
            content = root / "thesis.md"
            content.write_text("# 标题\n\n## 唯一章节\n\n- 唯一要点\n", encoding="utf-8")
            output = root / "filled.pptx"

            result = fill_template(template, content, output)

            slides = extract_pptx_slide_text(output)
            self.assertEqual(4, len(slides))  # cover + toc + 1 content + ending
            self.assertIn("唯一章节", slides[2][1])
            joined = "\n".join(text for _, text in slides)
            self.assertNotIn("点击输入", joined)
            self.assertEqual("clean", result.verdict)

    def test_overflow_reports_slide_number(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            template = root / "template.pptx"
            _build_template(template)
            content = root / "thesis.md"
            content.write_text(
                "# 标题\n\n## 第一章\n\n- 短要点\n\n## 第二章\n\n- " + "长内容" * 200 + "\n",
                encoding="utf-8",
            )
            output = root / "filled.pptx"

            result = fill_template(template, content, output)

            self.assertTrue(result.overflow)
            self.assertEqual([4], sorted({issue.slide for issue in result.overflow}))
            self.assertGreater(result.overflow[0].overshoot_pct, 15)
            self.assertTrue(result.verdict.startswith("needs-review"))
            report = result.report_path.read_text(encoding="utf-8")
            self.assertIn("## Operations", report)
            self.assertIn("## Overflow Issues", report)
            self.assertIn("## Stale Issues", report)
            self.assertIn("- slide 4:", report)
            self.assertIn("verdict: needs-review", report.splitlines()[-2] + report.splitlines()[-1])

    def test_stale_scan_flags_placeholders_and_failed_targets_only(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            template = root / "template.pptx"
            _build_template(template, planted_placeholder=True)
            content = root / "thesis.md"
            # No meta line -> cover subtitle is never targeted and must not be flagged.
            content.write_text("# 标题\n\n## 第一章\n\n- 一\n\n## 第二章\n\n- 二\n", encoding="utf-8")
            mapping = root / "map.json"
            mapping.write_text(json.dumps({"1": {"根本不存在的字符串": "替换值"}}), encoding="utf-8")
            output = root / "filled.pptx"

            result = fill_template(template, content, output, mapping_json=mapping)

            kinds = {(issue.kind, issue.slide) for issue in result.stale}
            self.assertIn(("placeholder", 4), kinds)  # planted 点击输入标题 survives
            self.assertIn(("target-not-found", 1), kinds)  # injected bogus target
            self.assertEqual(2, len(result.stale))
            for issue in result.stale:
                self.assertNotIn("张三", issue.text)
            slides = extract_pptx_slide_text(output)
            self.assertIn("张三", slides[0][1])  # untouched, unflagged template subtitle

    def test_cli_template_fill_smoke(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            template = root / "template.pptx"
            _build_template(template)
            content = root / "thesis.md"
            content.write_text(_FOUR_SECTION_MD, encoding="utf-8")
            output = root / "filled.pptx"

            rc = cli.main([
                "template-fill", str(template),
                "--content", str(content),
                "-o", str(output),
            ])

            self.assertEqual(0, rc)
            self.assertTrue(output.exists())
            self.assertTrue((root / "FILL-REPORT.md").exists())

    def test_cli_errors_on_sectionless_markdown(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            template = root / "template.pptx"
            _build_template(template)
            content = root / "empty.md"
            content.write_text("# 只有标题没有章节\n", encoding="utf-8")

            rc = cli.main([
                "template-fill", str(template),
                "--content", str(content),
                "-o", str(root / "out.pptx"),
            ])

            self.assertEqual(1, rc)

    def test_analyze_template_roles_and_slots(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            template = root / "template.pptx"
            _build_template(template)

            profile = analyze_template(template)

            self.assertEqual(
                ["cover", "toc", "content", "content", "ending"],
                [slide.role for slide in profile.slides],
            )
            self.assertEqual([3, 4], profile.content_numbers)
            self.assertEqual("论文题目样板", profile.slides[0].title.text)
            self.assertEqual("点击输入正文内容", profile.slides[2].body.text)


if __name__ == "__main__":
    unittest.main()
