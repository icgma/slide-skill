"""Wiring tests for v1.4 integration: enhancements, animations_v2, i18n, CLI, demo."""

from __future__ import annotations

import base64
import json
import sys
import unittest
from io import StringIO
from pathlib import Path
from tempfile import TemporaryDirectory
from xml.etree import ElementTree as ET

from slide_skill import cli, enhancements
from slide_skill.project import init_project
from slide_skill.svg_pipeline import create_spec, finalize_svg, generate_svg
from slide_skill.themes import get_theme


SAMPLE_SVG_TEMPLATE = """<svg xmlns="http://www.w3.org/2000/svg" width="1280" height="720" viewBox="0 0 1280 720">
  <g id="bg"><rect width="1280" height="720" fill="#0F172A"/></g>
  <g id="content">
  {placeholders}
  </g>
</svg>"""


def _b64(s: str) -> str:
    return base64.b64encode(s.encode("utf-8")).decode("ascii")


class EnhancementsTests(unittest.TestCase):
    def test_chart_placeholder_expands_to_g(self) -> None:
        spec = {"kind": "bar", "categories": ["A", "B"], "series": [{"name": "S", "values": [1, 2]}]}
        ph = (
            f'<g data-enhance="chart" data-spec-b64="{_b64(json.dumps(spec))}" '
            f'data-x="40" data-y="100" data-w="600" data-h="320"/>'
        )
        out = enhancements.expand_enhancements(SAMPLE_SVG_TEMPLATE.format(placeholders=ph))
        self.assertNotIn('data-enhance="chart"', out)
        self.assertIn("translate(40,100)", out)

    def test_code_placeholder_expands_with_text(self) -> None:
        code_b64 = _b64("print(1)\nprint(2)")
        ph = (
            f'<g data-enhance="code" data-language="python" '
            f'data-text-b64="{code_b64}" '
            f'data-x="0" data-y="0" data-w="800" data-line-numbers="true" data-highlight="1"/>'
        )
        out = enhancements.expand_enhancements(SAMPLE_SVG_TEMPLATE.format(placeholders=ph))
        self.assertNotIn('data-enhance="code"', out)
        # Code SVG always wraps in <g transform="translate(...)">
        self.assertIn("<g", out)

    def test_icon_placeholder_expands_with_size(self) -> None:
        ph = '<g data-enhance="icon" data-name="rocket" data-x="10" data-y="20" data-size="64"/>'
        out = enhancements.expand_enhancements(SAMPLE_SVG_TEMPLATE.format(placeholders=ph))
        self.assertNotIn('data-enhance="icon"', out)
        self.assertIn('width="64"', out)
        self.assertIn('x="10"', out)

    def test_bad_placeholder_becomes_comment_not_abort(self) -> None:
        ph = '<g data-enhance="chart" data-spec="not-json"/>'
        out = enhancements.expand_enhancements(SAMPLE_SVG_TEMPLATE.format(placeholders=ph))
        self.assertIn("<!-- enhancement chart skipped", out)

    def test_unknown_kind_left_alone(self) -> None:
        ph = '<g data-enhance="hologram"/>'
        text = SAMPLE_SVG_TEMPLATE.format(placeholders=ph)
        out = enhancements.expand_enhancements(text)
        self.assertEqual(text, out)

    def test_self_closing_and_empty_body_both_match(self) -> None:
        for body in ('<g data-enhance="icon" data-name="rocket" data-x="0" data-y="0"/>',
                     '<g data-enhance="icon" data-name="rocket" data-x="0" data-y="0"></g>'):
            out = enhancements.expand_enhancements(SAMPLE_SVG_TEMPLATE.format(placeholders=body))
            self.assertNotIn('data-enhance="icon"', out)


class FinalizeSvgIntegrationTests(unittest.TestCase):
    def test_finalize_expands_placeholders_and_uses_lock_theme(self) -> None:
        with TemporaryDirectory() as base:
            project = init_project("wiring", "ppt169", base, overwrite=True)
            src_md = project / "sources" / "input.md"
            src_md.write_text("# Title\n\n- a\n- b\n", encoding="utf-8")
            create_spec(project, src_md, theme_name="dark-tech")

            # Hand-author one SVG with an icon placeholder so we can verify
            # finalize_svg applies enhancements without depending on the
            # programmatic generator emitting placeholders.
            out = project / "svg_output"
            out.mkdir(exist_ok=True)
            (out / "slide_01.svg").write_text(
                SAMPLE_SVG_TEMPLATE.format(
                    placeholders='<g data-enhance="icon" data-name="rocket" data-x="50" data-y="50" data-size="32"/>'
                ),
                encoding="utf-8",
            )
            paths = finalize_svg(project)
            self.assertEqual(len(paths), 1)
            final_text = paths[0].read_text(encoding="utf-8")
            self.assertNotIn('data-enhance="icon"', final_text)
            self.assertIn('width="32"', final_text)


class SpecLockLangTests(unittest.TestCase):
    def test_create_spec_persists_detected_language(self) -> None:
        with TemporaryDirectory() as base:
            project = init_project("zh", "ppt169", base, overwrite=True)
            src_md = project / "sources" / "input.md"
            src_md.write_text("# 你好世界\n\n- 中文内容\n- 更多中文\n", encoding="utf-8")
            create_spec(project, src_md, theme_name="dark-tech")
            lock = json.loads((project / "spec_lock.json").read_text(encoding="utf-8"))
            self.assertEqual(lock.get("lang"), "zh")


class AnimationsV2WiringTests(unittest.TestCase):
    def test_v2_attrs_produce_p_timing(self) -> None:
        try:
            from pptx import Presentation  # noqa: F401
        except ImportError:
            self.skipTest("python-pptx not installed")
        from slide_skill.exporter import export_project

        with TemporaryDirectory() as base:
            project = init_project("anim2", "ppt169", base, overwrite=True)
            src_md = project / "sources" / "input.md"
            src_md.write_text("# Cover\n\n- bullet\n", encoding="utf-8")
            create_spec(project, src_md)
            # Hand-author SVG with v2 anim attrs on a real shape group.
            out = project / "svg_output"
            out.mkdir(exist_ok=True)
            svg = (
                '<svg xmlns="http://www.w3.org/2000/svg" width="1280" height="720" viewBox="0 0 1280 720">'
                '<g id="bg"><rect width="1280" height="720" fill="#0F172A"/></g>'
                '<g id="content" data-anim-v2-effect="fadeIn" data-anim-v2-trigger="onClick" '
                'data-anim-v2-duration="700" data-anim-v2-delay="200">'
                '<rect x="100" y="100" width="200" height="80" fill="#3B82F6"/>'
                '</g>'
                '</svg>'
            )
            (out / "slide_01.svg").write_text(svg, encoding="utf-8")
            finalize_svg(project)
            pptx = export_project(project)
            self.assertTrue(pptx.exists())
            import zipfile
            with zipfile.ZipFile(pptx) as zf:
                slide1 = zf.read("ppt/slides/slide1.xml").decode("utf-8")
            self.assertIn("p:timing", slide1)
            # v2 effect "fadeIn" → presetID 10
            self.assertIn('presetID="10"', slide1)


class CLIWiringTests(unittest.TestCase):
    def _run_cli(self, argv: list[str]) -> tuple[int, str, str]:
        out, err = StringIO(), StringIO()
        old_out, old_err = sys.stdout, sys.stderr
        sys.stdout, sys.stderr = out, err
        try:
            rc = cli.main(argv)
        finally:
            sys.stdout, sys.stderr = old_out, old_err
        return rc, out.getvalue(), err.getvalue()

    def test_html_preview_subcommand_writes_file(self) -> None:
        with TemporaryDirectory() as base:
            project = init_project("cli-html", "ppt169", base, overwrite=True)
            src_md = project / "sources" / "input.md"
            src_md.write_text("# Hi\n\n- a\n", encoding="utf-8")
            create_spec(project, src_md)
            generate_svg(project, src_md)
            finalize_svg(project)
            out_path = project / "exports" / "preview.html"
            rc, stdout, stderr = self._run_cli(["html-preview", str(project), "-o", str(out_path)])
            self.assertEqual(rc, 0, msg=stderr)
            self.assertTrue(out_path.exists())
            self.assertIn("<html", out_path.read_text(encoding="utf-8").lower())

    def test_font_preflight_subcommand_runs(self) -> None:
        with TemporaryDirectory() as base:
            project = init_project("cli-pf", "ppt169", base, overwrite=True)
            src_md = project / "sources" / "input.md"
            src_md.write_text("# Hello world\n", encoding="utf-8")
            create_spec(project, src_md)
            rc, stdout, stderr = self._run_cli(["font-preflight", str(project)])
            # rc may be 0 (no warn) or 1 (warn). We only care it doesn't crash.
            self.assertIn(rc, (0, 1), msg=stderr)
            self.assertIn("language:", stdout)


class DemoPreviewRouteTests(unittest.TestCase):
    def test_preview_html_route_serves_html(self) -> None:
        sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "tools" / "slide-demo"))
        try:
            import importlib
            app_mod = importlib.import_module("app")
        except Exception:  # noqa: BLE001
            self.skipTest("demo app not importable (likely missing flask)")
            return

        client = app_mod.app.test_client()

        # Bad job_id → 400
        self.assertEqual(client.get("/preview/!!bad!!/html").status_code, 400)

        # Build a real job dir with an init_project + finalize_svg.
        with TemporaryDirectory() as base:
            project = init_project("demo-job", "ppt169", base, overwrite=True)
            src_md = project / "sources" / "input.md"
            src_md.write_text("# Hi\n\n- a\n", encoding="utf-8")
            create_spec(project, src_md)
            generate_svg(project, src_md)
            finalize_svg(project)

            job_id = "abc123demo"
            job_root = app_mod.OUTPUT_ROOT / job_id
            if job_root.exists():
                import shutil as _sh
                _sh.rmtree(job_root)
            job_root.mkdir(parents=True)
            # Symlink/copy our project into the job_root so the route can find it.
            import shutil as _sh
            _sh.copytree(project, job_root / project.name)

            try:
                resp = client.get(f"/preview/{job_id}/html")
                self.assertEqual(resp.status_code, 200)
                self.assertIn(b"<html", resp.data.lower())
            finally:
                _sh.rmtree(job_root, ignore_errors=True)


if __name__ == "__main__":  # pragma: no cover
    unittest.main()
