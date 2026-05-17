from __future__ import annotations

import tempfile
import unittest
import zipfile
from pathlib import Path

from slide_skill.exporter import export_project, validate_pptx
from slide_skill.project import init_project
from slide_skill.svg_pipeline import create_spec, finalize_svg, generate_svg
from slide_skill.template_ops import delete_slides, duplicate_slide, inspect_template, reorder_slides, replace_text


class TemplateOpsTest(unittest.TestCase):
    def test_template_operations_preserve_valid_pptx(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source = root / "source.md"
            source.write_text("# Demo\n\n## One\n\nFirst.\n\n## Two\n\nSecond.\n", encoding="utf-8")
            project = init_project("Template Deck", base_dir=root / "projects")
            create_spec(project, source)
            generate_svg(project, source)
            finalize_svg(project)
            deck = export_project(project)

            info = inspect_template(deck)
            self.assertGreaterEqual(info["slide_count"], 2)

            replaced = root / "replaced.pptx"
            replace_text(deck, replaced, {"Demo": "Updated Demo"})
            self.assertTrue(validate_pptx(replaced)[0])

            duplicated = root / "duplicated.pptx"
            duplicate_slide(deck, duplicated, 1)
            self.assertTrue(validate_pptx(duplicated)[0])
            self.assertGreater(inspect_template(duplicated)["slide_count"], info["slide_count"])

            deleted = root / "deleted.pptx"
            delete_slides(deck, deleted, [1])
            self.assertTrue(validate_pptx(deleted)[0])
            deleted_info = inspect_template(deleted)
            self.assertEqual(deleted_info["slide_count"], info["slide_count"] - 1)

            reordered = root / "reordered.pptx"
            reorder_slides(deck, reordered, list(reversed(range(1, info["slide_count"] + 1))))
            self.assertTrue(validate_pptx(reordered)[0])
            reordered_info = inspect_template(reordered)
            self.assertIn("Two", reordered_info["slides"][0]["text"])

    def test_template_replace_handles_split_runs(self) -> None:
        from pptx import Presentation
        from pptx.util import Inches

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            deck = root / "split-runs.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
            paragraph = box.text_frame.paragraphs[0]
            paragraph.add_run().text = "Hello "
            paragraph.add_run().text = "World"
            prs.save(deck)

            output = root / "replaced.pptx"
            replace_text(deck, output, {"Hello World": "Updated"})

            info = inspect_template(output)
            self.assertIn("Updated", info["slides"][0]["text"])
            self.assertNotIn("Hello World", info["slides"][0]["text"])

    def test_delete_slide_removes_unreferenced_media(self) -> None:
        from PIL import Image
        from pptx import Presentation
        from pptx.util import Inches

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            image = root / "image.png"
            Image.new("RGB", (16, 16), "red").save(image)

            deck = root / "media.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(str(image), Inches(1), Inches(1), Inches(1), Inches(1))
            prs.save(deck)

            output = root / "deleted.pptx"
            delete_slides(deck, output, [1])

            with zipfile.ZipFile(output) as zf:
                self.assertEqual([], [name for name in zf.namelist() if name.startswith("ppt/media/")])


if __name__ == "__main__":
    unittest.main()
