"""SVG tag converter registry and converter functions — slide-skill v2.0."""

from __future__ import annotations

import math
import re
from typing import Callable
from xml.etree import ElementTree as ET

from .i18n import CJK_CHAR_WIDTH, LATIN_CHAR_WIDTH, detect_language, get_language_profile


def safe_float(val: str | None, default: float = 0.0) -> float:
    """Safely parse a float from an SVG attribute, ignoring units like px, em."""
    if not val:
        return default
    try:
        # Extract the first sequence that looks like a float (including negative and decimal)
        m = re.search(r"[-+]?[0-9]*\.?[0-9]+", str(val).strip())
        if m:
            return float(m.group(0))
        return default
    except ValueError:
        return default


def parse_hex(value: str | None) -> tuple[int, int, int] | None:
    if not value or value == "none":
        return None
    value = value.strip()
    if value.startswith("url("):
        return None  # gradient/pattern reference — handled separately
    value = value.lstrip("#")
    if len(value) != 6:
        return None
    try:
        return int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16)
    except ValueError:
        return None


def _parse_stop_color(stop_elem: ET.Element) -> tuple[int, int, int] | None:
    """Extract stop colour from a <stop> element, handling hex, rgb(), and style=""."""
    raw = stop_elem.attrib.get("stop-color", "")
    # Also check inline style for stop-color
    style = stop_elem.attrib.get("style", "")
    if not raw and style:
        m = re.search(r"stop-color\s*:\s*([^;]+)", style)
        if m:
            raw = m.group(1).strip()
    raw = raw.strip()
    if not raw or raw == "none":
        return None
    # Handle rgb(r, g, b)
    m_rgb = re.match(r"rgb\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*\)", raw, re.IGNORECASE)
    if m_rgb:
        return int(m_rgb.group(1)), int(m_rgb.group(2)), int(m_rgb.group(3))
    return parse_hex(raw)


def _parse_stop_offset(stop_elem: ET.Element) -> float:
    """Return stop offset in [0.0, 1.0]."""
    raw = stop_elem.attrib.get("offset", "0")
    # Also check inline style
    style = stop_elem.attrib.get("style", "")
    if style:
        m = re.search(r"\boffset\s*:\s*([^;]+)", style)
        if m:
            raw = m.group(1).strip()
    raw = raw.strip()
    if raw.endswith("%"):
        try:
            return float(raw[:-1]) / 100.0
        except ValueError:
            return 0.0
    return safe_float(raw)


def local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1]


def _find_gradient_elem(root: ET.Element, ref_id: str) -> ET.Element | None:
    """Locate a gradient element by id anywhere in <defs>."""
    for defs_elem in root.iter():
        if local_name(defs_elem.tag) != "defs":
            continue
        for grad in defs_elem:
            grad_tag = local_name(grad.tag)
            if grad_tag not in {"linearGradient", "radialGradient"}:
                continue
            if grad.attrib.get("id") == ref_id:
                return grad
    return None


def _resolve_xlink_href(grad: ET.Element, root: ET.Element) -> ET.Element:
    """Follow xlink:href / href to get the element that actually holds the stops."""
    href = grad.attrib.get("{http://www.w3.org/1999/xlink}href") or grad.attrib.get("href", "")
    if href.startswith("#"):
        target = _find_gradient_elem(root, href[1:])
        if target is not None:
            return target
    return grad


def extract_gradient_info(root: ET.Element, ref_id: str) -> dict | None:
    """
    Return a dict with keys:
      type  - "linear" | "radial"
      stops - list of (pos_0_1, (r, g, b))  sorted by pos
      ang   - DrawingML angle in 1/60000 of a degree (linear only)
    Returns None if gradient not found or has no usable stops.
    """
    grad = _find_gradient_elem(root, ref_id)
    if grad is None:
        return None
    grad_tag = local_name(grad.tag)

    stop_source = _resolve_xlink_href(grad, root)
    raw_stops = [s for s in stop_source if local_name(s.tag) == "stop"]
    if not raw_stops:
        raw_stops = [s for s in grad if local_name(s.tag) == "stop"]
    if not raw_stops:
        return None

    stops: list[tuple[float, tuple[int, int, int]]] = []
    for s in raw_stops:
        color = _parse_stop_color(s)
        if color is None:
            continue
        pos = max(0.0, min(1.0, _parse_stop_offset(s)))
        stops.append((pos, color))
    stops.sort(key=lambda t: t[0])
    if not stops:
        return None

    if grad_tag == "linearGradient":
        def _pct(v: str) -> float:
            v = v.strip()
            if v.endswith("%"):
                try:
                    return float(v[:-1]) / 100.0
                except ValueError:
                    return 0.0
            return safe_float(v)

        x1 = _pct(grad.attrib.get("x1", "0"))
        y1 = _pct(grad.attrib.get("y1", "0"))
        x2 = _pct(grad.attrib.get("x2", "1"))
        y2 = _pct(grad.attrib.get("y2", "0"))
        dx = x2 - x1
        dy = y2 - y1
        angle_rad = math.atan2(dx, dy)
        angle_deg = math.degrees(angle_rad) % 360
        ang = int(round(angle_deg * 60000))
        return {"type": "linear", "stops": stops, "ang": ang}

    return {"type": "radial", "stops": stops, "ang": 0}


def extract_gradient_color(root: ET.Element, ref_id: str) -> tuple[int, int, int] | None:
    """Return the midpoint stop colour of a gradient referenced by id (without #).
    Kept for backward-compatibility; prefer extract_gradient_info for native rendering."""
    info = extract_gradient_info(root, ref_id)
    if info and info["stops"]:
        stops = info["stops"]
        mid = stops[len(stops) // 2]
        return mid[1]
    return None


def resolve_fill(elem: ET.Element, root: ET.Element) -> tuple[int, int, int] | None:
    """Resolve fill attribute to a solid colour only (gradient refs → None)."""
    fill = elem.attrib.get("fill", "")
    if fill.startswith("url(#"):
        return None  # caller should use _apply_native_gradient instead
    return parse_hex(fill)


def _is_gradient_fill(elem: ET.Element) -> tuple[bool, str]:
    """Return (True, ref_id) when fill is a gradient url reference."""
    fill = elem.attrib.get("fill", "")
    if fill.startswith("url(#"):
        return True, fill[5:].rstrip(")")
    return False, ""


def _build_grad_fill_xml_string(info: dict) -> str:
    """Build a DrawingML <a:gradFill> XML string from gradient info."""
    stops_xml = ""
    for pos, (r, g, b) in info["stops"]:
        pos_val = int(round(pos * 100000))
        hex_val = f"{r:02X}{g:02X}{b:02X}"
        stops_xml += f'<a:gs pos="{pos_val}"><a:srgbClr val="{hex_val}"/></a:gs>'

    if info["type"] == "linear":
        direction_xml = f'<a:lin ang="{info["ang"]}" scaled="0"/>'
    else:
        direction_xml = (
            '<a:path path="circle">'
            '<a:fillToRect l="50000" t="50000" r="50000" b="50000"/>'
            "</a:path>"
        )

    return (
        '<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f"<a:gsLst>{stops_xml}</a:gsLst>"
        f"{direction_xml}"
        "</a:gradFill>"
    )


def _apply_native_gradient(shape, grad_info: dict) -> None:
    """Replace spPr fill with a native DrawingML gradFill element using lxml."""
    try:
        from pptx.oxml import parse_xml
    except ImportError:
        return

    spPr = shape._element.spPr  # type: ignore[attr-defined]

    for child in list(spPr):
        tag_local = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag_local in ("solidFill", "gradFill", "noFill", "blipFill", "pattFill", "grpFill"):
            spPr.remove(child)

    xml_str = _build_grad_fill_xml_string(grad_info)
    grad_fill_elem = parse_xml(xml_str)

    # In OOXML spPr ordering, fill must come AFTER geometry (prstGeom/custGeom)
    # and BEFORE line/effect children.  Scan all children to find the last
    # pre-fill element so we always insert in a schema-valid position.
    _PRE_FILL_TAGS = {"xfrm", "prstGeom", "custGeom"}
    insert_pos = 0
    for i, child in enumerate(spPr):
        tag_local = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag_local in _PRE_FILL_TAGS:
            insert_pos = i + 1  # keep updating — insert after the LAST pre-fill element

    spPr.insert(insert_pos, grad_fill_elem)


def _apply_fill_and_line_with_gradient(shape, elem: ET.Element, rgb_cls, root: ET.Element | None, grad_info: dict) -> None:
    """Apply native gradient fill and then line styles."""
    _apply_native_gradient(shape, grad_info)
    _apply_line(shape.line, elem, rgb_cls, root)


class ConverterRegistry:
    """Extensible dispatch for SVG element → PPTX shape converters."""

    def __init__(self) -> None:
        self._converters: dict[str, Callable] = {}
        self._root: ET.Element | None = None

    def register(self, tag: str, converter: Callable) -> None:
        self._converters[tag] = converter

    def set_root(self, root: ET.Element) -> None:
        """Store SVG root for gradient/defs resolution."""
        self._root = root

    def dispatch(self, slide, elem: ET.Element, scale_x: float, scale_y: float, meta: dict, rgb_cls) -> None:
        tag = local_name(elem.tag)
        converter = self._converters.get(tag)
        if converter:
            converter(slide, elem, scale_x, scale_y, meta, rgb_cls, self._root)

    def supported_tags(self) -> set[str]:
        return set(self._converters.keys())


def _box(elem, scale_x, scale_y, width_attr="width", height_attr="height"):
    x = safe_float(elem.attrib.get("x", 0)) * scale_x
    y = safe_float(elem.attrib.get("y", 0)) * scale_y
    w = safe_float(elem.attrib.get(width_attr, 0)) * scale_x
    h = safe_float(elem.attrib.get(height_attr, 0)) * scale_y
    return x, y, max(w, 0.01), max(h, 0.01)


def _oval_box(elem, scale_x, scale_y):
    tag = local_name(elem.tag)
    cx = safe_float(elem.attrib.get("cx", 0))
    cy = safe_float(elem.attrib.get("cy", 0))
    if tag == "circle":
        r = safe_float(elem.attrib.get("r", 0))
        return (cx - r) * scale_x, (cy - r) * scale_y, (2 * r) * scale_x, (2 * r) * scale_y
    rx = safe_float(elem.attrib.get("rx", 0))
    ry = safe_float(elem.attrib.get("ry", 0))
    return (cx - rx) * scale_x, (cy - ry) * scale_y, (2 * rx) * scale_x, (2 * ry) * scale_y


def _apply_fill_and_line(shape, elem, rgb_cls, root=None):
    is_grad, ref_id = _is_gradient_fill(elem)
    if is_grad and root is not None:
        grad_info = extract_gradient_info(root, ref_id)
        if grad_info:
            _apply_fill_and_line_with_gradient(shape, elem, rgb_cls, root, grad_info)
            return

    fill = resolve_fill(elem, root) if root is not None else parse_hex(elem.attrib.get("fill"))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb_cls(*fill)
        # Apply opacity if present
        opacity_str = elem.attrib.get("opacity") or elem.attrib.get("fill-opacity")
        if opacity_str:
            opacity = safe_float(opacity_str, default=-1.0)
            if 0 <= opacity < 1:
                try:
                    from pptx.oxml.ns import qn
                    from lxml import etree as _lxml_etree
                    spPr = shape._element.spPr
                    solid = spPr.find(qn("a:solidFill"))
                    if solid is not None:
                        srgb = solid.find(qn("a:srgbClr"))
                        if srgb is not None:
                            alpha_el = _lxml_etree.SubElement(srgb, qn("a:alpha"))
                            alpha_el.set("val", str(int(opacity * 100000)))
                except (ValueError, AttributeError):
                    pass  # Opacity application is best-effort
    else:
        shape.fill.background()
    _apply_line(shape.line, elem, rgb_cls, root)


def _apply_line(line, elem, rgb_cls, root=None):
    stroke = parse_hex(elem.attrib.get("stroke"))
    if stroke:
        line.color.rgb = rgb_cls(*stroke)
        if elem.attrib.get("stroke-width"):
            line.width = int(safe_float(elem.attrib["stroke-width"]) * 12700)
    else:
        line.fill.background()


def _parse_simple_translate(transform: str) -> tuple[float, float]:
    """Extract (tx, ty) from translate(tx, ty) or translate(tx). Returns (0, 0) if not parseable."""
    m = re.match(r"translate\(\s*([\d.\-]+)(?:[,\s]+([\d.\-]+))?\s*\)", transform.strip())
    if m:
        tx = safe_float(m.group(1))
        ty = safe_float(m.group(2)) if m.group(2) else 0.0
        return tx, ty
    return 0.0, 0.0


# --- converter functions ---

def convert_rect(slide, elem, scale_x, scale_y, meta, rgb_cls, root=None):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    x, y, w, h = _box(elem, scale_x, scale_y)

    # Apply simple translate transform
    transform = elem.attrib.get("transform", "")
    if "translate" in transform:
        tx, ty = _parse_simple_translate(transform)
        x += tx * scale_x
        y += ty * scale_y

    rx = safe_float(elem.attrib.get("rx", "0"))
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rx else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    _apply_fill_and_line(shape, elem, rgb_cls, root)


def convert_oval(slide, elem, scale_x, scale_y, meta, rgb_cls, root=None):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    x, y, w, h = _oval_box(elem, scale_x, scale_y)
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    _apply_fill_and_line(shape, elem, rgb_cls, root)


def convert_line(slide, elem, scale_x, scale_y, meta, rgb_cls, root=None):
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.util import Inches

    x1 = safe_float(elem.attrib.get("x1", 0)) * scale_x
    y1 = safe_float(elem.attrib.get("y1", 0)) * scale_y
    x2 = safe_float(elem.attrib.get("x2", 0)) * scale_x
    y2 = safe_float(elem.attrib.get("y2", 0)) * scale_y
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    _apply_line(connector.line, elem, rgb_cls, root)


def convert_text(slide, elem, scale_x, scale_y, meta, rgb_cls, root=None):
    from pptx.util import Inches, Pt

    tspans = [c for c in elem if local_name(c.tag) == "tspan"]
    if tspans:
        para_texts = [(c.text or "").strip() for c in tspans]
        para_texts = [t for t in para_texts if t]
        text = "\n".join(para_texts)
    else:
        para_texts = []
        text = "".join(elem.itertext()).strip()
    if not text:
        return
    x = safe_float(elem.attrib.get("x", 0)) * scale_x
    font_size_svg = safe_float(elem.attrib.get("font-size", "24"), default=24.0)
    font_pt = font_size_svg * scale_y * 72.0
    baseline = safe_float(elem.attrib.get("y", 0)) * scale_y
    ascent_in = (font_pt / 72.0) * 0.92
    y = max(0.0, baseline - ascent_in)

    anchor = elem.attrib.get("text-anchor", "start").lower()
    if para_texts:
        w = max(0.5, max(approx_w_in(t, font_pt) for t in para_texts))
    elif anchor in ("middle", "end"):
        w = max(0.5, approx_w_in(text, font_pt))
    else:
        w = max(0.5, float(meta["canvas"]["pptx_width_in"]) - x - 0.5)
    if anchor == "middle":
        x = max(0.0, x - w / 2.0)
    elif anchor == "end":
        x = max(0.0, x - w)

    bold = elem.attrib.get("font-weight", "").lower() in {"bold", "700", "800", "900"}
    fill_val = elem.attrib.get("fill", "#111827")
    color = parse_hex(fill_val)
    family = elem.attrib.get("font-family")
    primary = None
    cjk = None
    if family:
        families = [f.strip().strip("'").strip('"') for f in family.split(",")]
        families = [f for f in families if f and f.lower() not in {"sans-serif", "serif", "monospace", "cursive", "fantasy"}]
        cjk_markers = ("yahei", "pingfang", "noto sans cjk", "noto serif cjk", "noto sans sc", "noto serif sc",
                       "source han", "songti", "stsong", "stheiti", "simsun", "simhei", "fangsong", "kaiti",
                       "hiragino", "yu gothic", "meiryo", "ms gothic", "ms mincho",
                       "malgun gothic", "nanum", "apple sd gothic")
        primary = families[0] if families else family
        cjk = next((f for f in families if any(m in f.lower() for m in cjk_markers)), None)

    _create_styled_textbox(
        slide, text, x, y, w, font_pt, bold, color, primary, rgb_cls,
        cjk=cjk, anchor=anchor, word_wrap=not para_texts)


def convert_image(slide, elem, scale_x, scale_y, meta, rgb_cls, root=None):
    from pptx.util import Inches
    from pathlib import Path

    href = elem.attrib.get("href") or elem.attrib.get("{http://www.w3.org/1999/xlink}href")
    if not href:
        return
    image_path = (Path(".") / href).resolve()
    if not image_path.exists():
        return
    x, y, w, h = _box(elem, scale_x, scale_y, width_attr="width", height_attr="height")
    slide.shapes.add_picture(str(image_path), Inches(x), Inches(y), Inches(w), Inches(h))


def convert_path(slide, elem, scale_x, scale_y, meta, rgb_cls, root=None):
    from .geometry import build_freeform_xml, compute_bbox, parse_svg_path

    d = elem.attrib.get("d", "")
    if not d.strip():
        return
    commands = parse_svg_path(d)
    _add_freeform(slide, elem, commands, scale_x, scale_y, rgb_cls, root)


def convert_polygon(slide, elem, scale_x, scale_y, meta, rgb_cls, root=None):
    from .geometry import parse_svg_points, points_to_commands

    pts = parse_svg_points(elem.attrib.get("points", ""))
    commands = points_to_commands(pts, closed=True)
    _add_freeform(slide, elem, commands, scale_x, scale_y, rgb_cls, root)


def convert_polyline(slide, elem, scale_x, scale_y, meta, rgb_cls, root=None):
    from .geometry import parse_svg_points, points_to_commands

    pts = parse_svg_points(elem.attrib.get("points", ""))
    commands = points_to_commands(pts, closed=False)
    _add_freeform(slide, elem, commands, scale_x, scale_y, rgb_cls, root)


def _add_freeform(slide, elem, commands, scale_x, scale_y, rgb_cls, root=None):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches
    from .geometry import EMU_PER_INCH, build_freeform_xml, compute_bbox

    if not commands:
        return

    min_x, min_y, max_x, max_y = compute_bbox(commands)
    bbox_w = max(max_x - min_x, 0.01)
    bbox_h = max(max_y - min_y, 0.01)
    left_in = min_x * scale_x
    top_in = min_y * scale_y
    width_in = bbox_w * scale_x
    height_in = bbox_h * scale_y
    w_emu = int(width_in * EMU_PER_INCH)
    h_emu = int(height_in * EMU_PER_INCH)

    if w_emu <= 0 or h_emu <= 0:
        return

    cust_geom = build_freeform_xml(commands, w_emu, h_emu, min_x, min_y, scale_x, scale_y)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left_in), Inches(top_in),
        Inches(width_in), Inches(height_in),
    )
    spPr = shape._element.spPr
    for child in list(spPr):
        tag_local = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag_local in ("prstGeom", "custGeom"):
            spPr.remove(child)
    spPr.insert(0, cust_geom)
    _apply_fill_and_line(shape, elem, rgb_cls, root)


def _noop_converter(slide, elem, scale_x, scale_y, meta, rgb_cls, root=None):
    """Silently skip non-drawable structural tags (defs, g, filter, etc.)."""
    pass


def approx_w_in(s: str, font_pt: float) -> float:
    ems = sum(CJK_CHAR_WIDTH if ord(ch) >= 0x2E80 else LATIN_CHAR_WIDTH for ch in s)
    return (ems + 1.0) * (font_pt / 72.0)


def _create_styled_textbox(slide, text, x, y, w, font_pt, bold, color, family,
                           rgb_cls, cjk=None, anchor="start", word_wrap=True):
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    line_count = max(1, text.count("\n") + 1)
    h = max(0.35, (font_pt / 72.0) * 1.25 * line_count)
    textbox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    textbox.text = text
    frame = textbox.text_frame
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
    frame.word_wrap = word_wrap

    pp_align = None
    if anchor == "middle":
        pp_align = PP_ALIGN.CENTER
    elif anchor == "end":
        pp_align = PP_ALIGN.RIGHT

    for paragraph in frame.paragraphs:
        if pp_align is not None:
            paragraph.alignment = pp_align
        if not paragraph.runs:
            paragraph.add_run()
        for run in paragraph.runs:
            run.font.size = Pt(font_pt)
            run.font.bold = bold
            if color:
                run.font.color.rgb = rgb_cls(*color)
            if family:
                run.font.name = family
            if cjk:
                from pptx.oxml.ns import qn
                from lxml import etree as _lxml_etree
                rPr = run._r.get_or_add_rPr()
                for ea in rPr.findall(qn("a:ea")):
                    rPr.remove(ea)
                ea = _lxml_etree.SubElement(rPr, qn("a:ea"))
                ea.set("typeface", cjk)

    return textbox


def convert_bilingual_group(slide, group_elem, scale_x, scale_y, meta, rgb_cls, root=None):
    text_elems = [c for c in group_elem if local_name(c.tag) == "text"]
    if len(text_elems) < 2:
        return []

    zh_elems, en_elems = [], []
    for t in text_elems:
        content = "".join(t.itertext()).strip()
        if not content:
            continue
        lang = detect_language(content)
        if lang in ("zh", "ja", "ko"):
            zh_elems.append(t)
        else:
            en_elems.append(t)
    if not zh_elems or not en_elems:
        return []

    layout = group_elem.attrib.get("data-layout", "stacked")
    lang_size_attr = group_elem.attrib.get("data-lang-size", "")

    first_zh = zh_elems[0]
    first_en = en_elems[0]
    zh_font_svg = safe_float(first_zh.attrib.get("font-size", "24"), default=24.0)
    zh_font_pt_default = zh_font_svg * scale_y * 72.0
    zh_font_pt, en_font_pt = _parse_bilingual_font_sizes(lang_size_attr, zh_font_pt_default)

    base_x = safe_float(first_zh.attrib.get("x", 0)) * scale_x
    base_y = safe_float(first_zh.attrib.get("y", 0)) * scale_y
    canvas_w = float(meta["canvas"]["pptx_width_in"])

    zh_content = "\n".join("".join(t.itertext()).strip() for t in zh_elems)
    en_content = "\n".join("".join(t.itertext()).strip() for t in en_elems)

    zh_bold = first_zh.attrib.get("font-weight", "").lower() in {"bold", "700", "800", "900"}
    en_bold = first_en.attrib.get("font-weight", "").lower() in {"bold", "700", "800", "900"}
    zh_color = parse_hex(first_zh.attrib.get("fill", "#111827"))
    en_color = parse_hex(first_en.attrib.get("fill", "#111827"))
    zh_family = first_zh.attrib.get("font-family")
    en_family = first_en.attrib.get("font-family")

    shapes = []
    if layout == "side-by-side":
        total_w = max(1.0, canvas_w - base_x - 0.5)
        half_w = total_w * 0.48
        gap = total_w * 0.04
        zh_ascent = (zh_font_pt / 72.0) * 0.92
        y = max(0.0, base_y - zh_ascent)
        shapes.append(_create_styled_textbox(
            slide, zh_content, base_x, y, half_w, zh_font_pt,
            zh_bold, zh_color, zh_family, rgb_cls))
        shapes.append(_create_styled_textbox(
            slide, en_content, base_x + half_w + gap, y, half_w, en_font_pt,
            en_bold, en_color, en_family, rgb_cls))
    else:
        zh_w = max(0.5, approx_w_in(zh_content, zh_font_pt))
        en_w = max(0.5, approx_w_in(en_content, en_font_pt))
        w = max(zh_w, en_w)
        zh_ascent = (zh_font_pt / 72.0) * 0.92
        zh_y = max(0.0, base_y - zh_ascent)
        zh_line_count = max(1, zh_content.count("\n") + 1)
        zh_h = max(0.35, (zh_font_pt / 72.0) * 1.25 * zh_line_count)
        shapes.append(_create_styled_textbox(
            slide, zh_content, base_x, zh_y, w, zh_font_pt,
            zh_bold, zh_color, zh_family, rgb_cls))
        en_y = zh_y + zh_h + 0.05
        shapes.append(_create_styled_textbox(
            slide, en_content, base_x, en_y, w, en_font_pt,
            en_bold, en_color, en_family, rgb_cls))

    return shapes


def _parse_bilingual_font_sizes(attr, zh_font_pt):
    default_ratio = 0.75
    zh_size = zh_font_pt
    en_size = zh_font_pt * default_ratio
    if not attr:
        return zh_size, en_size
    if ";" in attr:
        for part in attr.split(";"):
            if ":" in part:
                key, val = part.split(":", 1)
                try:
                    size = float(val.rstrip("pt"))
                except ValueError:
                    continue
                if key.strip().lower() in ("zh", "ja", "ko"):
                    zh_size = size
                elif key.strip().lower() in ("en", "latin"):
                    en_size = size
    else:
        try:
            en_size = float(attr.rstrip("pt"))
        except ValueError:
            pass
    return zh_size, en_size


def create_default_registry() -> ConverterRegistry:
    """Create a registry with all built-in SVG tag converters."""
    reg = ConverterRegistry()
    reg.register("rect", convert_rect)
    reg.register("circle", convert_oval)
    reg.register("ellipse", convert_oval)
    reg.register("line", convert_line)
    reg.register("text", convert_text)
    reg.register("image", convert_image)
    reg.register("path", convert_path)
    reg.register("polygon", convert_polygon)
    reg.register("polyline", convert_polyline)
    # Silently skip structural/non-drawable tags
    for structural_tag in ("defs", "linearGradient", "radialGradient", "stop",
                           "filter", "feGaussianBlur", "feDropShadow", "feOffset", "feFlood",
                           "feComposite", "feMerge", "feMergeNode",
                           "clipPath", "mask", "pattern", "use", "title", "desc", "tspan", "g"):
        reg.register(structural_tag, _noop_converter)
    return reg
