"""SVG-as-image preview PPTX export for cross-platform sharing."""

from __future__ import annotations

import tempfile
from pathlib import Path

from .project import load_project
from .util import ensure_dir, timestamp


def export_preview_pptx(project_path: Path | str, output: Path | str | None = None, stage: str = "final") -> Path:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Emu
    except ImportError as exc:
        raise RuntimeError("PPTX export requires python-pptx. Install with: python -m pip install -e .") from exc

    project = Path(project_path)
    meta = load_project(project)
    svg_dir = project / ("svg_final" if stage == "final" else "svg_output")
    svg_files = sorted(svg_dir.glob("*.svg"))
    if not svg_files:
        raise FileNotFoundError(f"No SVG files found in {svg_dir}")

    prs = Presentation()
    prs.slide_width = Inches(float(meta["canvas"]["pptx_width_in"]))
    prs.slide_height = Inches(float(meta["canvas"]["pptx_height_in"]))
    blank_layout = prs.slide_layouts[6]

    img_dir = ensure_dir(project / "images" / "_preview_cache")
    rendered_images = _render_svgs_to_images(svg_files, img_dir, meta)

    for slide_index, (svg_file, img_path) in enumerate(zip(svg_files, rendered_images), start=1):
        slide = prs.slides.add_slide(blank_layout)
        if img_path and img_path.exists():
            left = Emu(0)
            top = Emu(0)
            slide.shapes.add_picture(str(img_path), left, top, prs.slide_width, prs.slide_height)
        else:
            _add_placeholder_text(slide, svg_file, meta)

    if output:
        out_path = Path(output)
    else:
        ts = timestamp()
        exports_dir = ensure_dir(project / "backup" / ts)
        out_path = exports_dir / f"{meta['name']}_preview.pptx"
    prs.save(str(out_path))
    return out_path


def _render_svgs_to_images(svg_files: list[Path], img_dir: Path, meta: dict) -> list[Path | None]:
    results: list[Path | None] = []
    width = int(meta["canvas"]["width"])
    height = int(meta["canvas"]["height"])
    for svg_file in svg_files:
        png_path = img_dir / (svg_file.stem + ".png")
        if png_path.exists():
            results.append(png_path)
            continue
        rendered = _try_render_svg(svg_file, png_path, width, height)
        results.append(rendered)
    return results


def _try_render_svg(svg_file: Path, png_path: Path, width: int, height: int) -> Path | None:
    try:
        import cairosvg
        cairosvg.svg2png(url=str(svg_file), write_to=str(png_path), output_width=width, output_height=height)
        return png_path
    except ImportError:
        pass
    except Exception:
        pass
    try:
        from PIL import Image
        img = Image.open(svg_file)
        img.save(png_path, "PNG")
        return png_path
    except Exception:
        pass
    return None


def _add_placeholder_text(slide, svg_file: Path, meta: dict) -> None:
    try:
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"Preview: {svg_file.name} (render not available)"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    except Exception:
        pass
