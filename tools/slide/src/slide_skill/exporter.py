"""SVG to native PPTX export — slide-skill v2.0."""

from __future__ import annotations

import re
import shutil
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from .converters import create_default_registry, local_name, convert_bilingual_group
from .clip_path import collect_clip_paths, resolve_clip_path, apply_clip_path_to_shape
from .filter_effects import collect_filters, resolve_filter, apply_filter_to_shape
from .gradient_fills import collect_gradients
from .intake import extract_pptx_slide_text
from .pattern_fill import collect_patterns, resolve_pattern_fill, apply_pattern_to_shape
from .project import load_project
from .svg_pipeline import check_project_svg
from .util import copy_tree, ensure_dir, timestamp

from .animations import inject_transition, inject_timing
from .animations_v2 import (
    EFFECT_CATALOG,
    VALID_TRIGGERS,
    ElementAnimation,
    inject_timing_v2,
)


def _safe_int(val: str | None, default: int = 0) -> int:
    """Safely parse an int from an SVG data attribute, stripping units like ms/s."""
    if not val:
        return default
    import re
    m = re.search(r"[-+]?\d+", str(val).strip())
    return int(m.group(0)) if m else default


def export_project(project_path: Path | str, output: Path | str | None = None, stage: str = "final", preview: bool = True) -> Path:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches
    except ImportError as exc:
        raise RuntimeError("PPTX export requires python-pptx. Install with: python -m pip install -e .") from exc

    project = Path(project_path)
    meta = load_project(project)
    svg_dir = project / ("svg_final" if stage == "final" else "svg_output")
    svg_files = sorted(svg_dir.glob("*.svg"))
    if not svg_files:
        raise FileNotFoundError(f"No SVG files found in {svg_dir}")
    ok_svg, svg_issues = check_project_svg(project, stage=stage)
    if not ok_svg:
        details = "\n".join(f"{issue.level}: {issue.file}: {issue.message}" for issue in svg_issues)
        raise RuntimeError(f"SVG quality gate failed before export:\n{details}")

    prs = Presentation()
    prs.slide_width = Inches(float(meta["canvas"]["pptx_width_in"]))
    prs.slide_height = Inches(float(meta["canvas"]["pptx_height_in"]))
    blank_layout = prs.slide_layouts[6]
    scale_x = float(meta["canvas"]["pptx_width_in"]) / float(meta["canvas"]["width"])
    scale_y = float(meta["canvas"]["pptx_height_in"]) / float(meta["canvas"]["height"])
    registry = create_default_registry()
    notes = _read_project_notes(project, len(svg_files))

    for slide_index, svg_file in enumerate(svg_files, start=1):
        slide = prs.slides.add_slide(blank_layout)
        root = ET.fromstring(svg_file.read_text(encoding="utf-8"))

        # Give the registry access to the SVG root for gradient resolution
        registry.set_root(root)

        slide_transition = None
        animated_groups: list[dict] = []
        animated_groups_v2: list[ElementAnimation] = []

        for group_elem in root:
            if local_name(group_elem.tag) != "g":
                continue
            transition = group_elem.attrib.get("data-transition")
            if transition:
                slide_transition = transition

            v2_effect = group_elem.attrib.get("data-anim-v2-effect")
            if v2_effect and v2_effect in EFFECT_CATALOG:
                trigger = group_elem.attrib.get("data-anim-v2-trigger", "onClick")
                if trigger not in VALID_TRIGGERS:
                    trigger = "onClick"
                try:
                    order_attr = group_elem.attrib.get("data-anim-v2-order")
                    animated_groups_v2.append(ElementAnimation(
                        target=str(len(animated_groups_v2)),
                        effect=v2_effect,
                        trigger=trigger,
                        order=_safe_int(order_attr, len(animated_groups_v2)) if order_attr is not None else len(animated_groups_v2),
                        delay=_safe_int(group_elem.attrib.get("data-anim-v2-delay", "0")),
                        duration=_safe_int(group_elem.attrib.get("data-anim-v2-duration", "500"), 500),
                    ))
                except (ValueError, KeyError):
                    pass
                continue  # v2 supersedes v1 for this group

            anim_type = group_elem.attrib.get("data-anim")
            if anim_type:
                animated_groups.append({
                    "type": anim_type,
                    "duration": _safe_int(group_elem.attrib.get("data-anim-duration", "500"), 500),
                    "delay": _safe_int(group_elem.attrib.get("data-anim-delay", "0")),
                })

        shape_count_before = len(slide.shapes)
        gradients = collect_gradients(root)
        clips = collect_clip_paths(root)
        patterns = collect_patterns(root)
        svg_filters = collect_filters(root)

        # Pre-scan for bilingual groups — their text children are handled as a unit
        bilingual_group_ids = set()
        bilingual_child_ids = set()
        for elem in root.iter():
            if local_name(elem.tag) == "g" and elem.attrib.get("data-bilingual") == "true":
                bilingual_group_ids.add(id(elem))
                for child in elem:
                    if local_name(child.tag) == "text":
                        bilingual_child_ids.add(id(child))

        drawable_elems = []
        for elem in root.iter():
            tag = local_name(elem.tag)

            # Skip children of bilingual groups — they're handled by the group
            if id(elem) in bilingual_child_ids:
                continue

            shape_count_pre = len(slide.shapes)

            if id(elem) in bilingual_group_ids:
                bilingual_shapes = convert_bilingual_group(
                    slide, elem, scale_x, scale_y, meta, RGBColor, root)
                if bilingual_shapes:
                    drawable_elems.append(elem)
            else:
                registry.dispatch(slide, elem, scale_x, scale_y, meta, RGBColor)
                if len(slide.shapes) > shape_count_pre and tag not in ("defs", "linearGradient", "radialGradient", "stop", "clipPath", "mask", "pattern", "filter", "feGaussianBlur", "feOffset", "feFlood", "feComposite", "feMerge", "feMergeNode", "style", "title", "desc", "metadata"):
                    drawable_elems.append(elem)
        shape_count_after = len(slide.shapes)

        new_shapes = list(slide.shapes)[shape_count_before:]

        parent_map = {c: p for p in root.iter() for c in p}

        for i, shape in enumerate(new_shapes):
            if i >= len(drawable_elems):
                break
            elem = drawable_elems[i]
            filter_attr = elem.attrib.get("filter", "")
            ancestor = parent_map.get(elem)
            while not filter_attr and ancestor is not None:
                filter_attr = ancestor.attrib.get("filter", "")
                ancestor = parent_map.get(ancestor)
            if filter_attr.startswith("url(#") and svg_filters:
                filt = resolve_filter(filter_attr, svg_filters)
                if filt:
                    apply_filter_to_shape(shape, filt, scale_x, scale_y)
            clip_attr = elem.attrib.get("clip-path") or elem.attrib.get("mask")
            if clip_attr and clips:
                clip = resolve_clip_path(clip_attr, clips)
                if clip and clip.get("commands"):
                    apply_clip_path_to_shape(shape, clip, scale_x, scale_y)
            fill_attr = elem.attrib.get("fill", "")
            if fill_attr.startswith("url(#") and patterns:
                pat = resolve_pattern_fill(fill_attr, patterns)
                if pat:
                    apply_pattern_to_shape(shape, pat, RGBColor)

        if slide_transition:
            inject_transition(slide._element, slide_transition)

        if animated_groups_v2:
            # Resolve target index -> sp_id from the just-emitted shapes.
            for i, anim in enumerate(animated_groups_v2):
                if i < len(new_shapes):
                    sp = new_shapes[i]._element
                    anim.sp_id = sp.get("id") or str(2 + i)
            inject_timing_v2(slide._element, animated_groups_v2)
        elif animated_groups:
            anim_entries = []
            for i, group_info in enumerate(animated_groups):
                if i < len(new_shapes):
                    sp = new_shapes[i]._element
                    sp_id = sp.get("id") or str(2 + i)
                    anim_entries.append({
                        "sp_id": sp_id,
                        "preset": group_info["type"],
                        "duration": group_info["duration"],
                        "delay": group_info["delay"],
                    })
            if anim_entries:
                inject_timing(slide._element, anim_entries)

        _embed_slide_notes(slide, notes[slide_index - 1])

    out_dir = ensure_dir(project / "exports")
    out_path = Path(output) if output else out_dir / f"{meta['name']}_{timestamp()}.pptx"
    prs.save(out_path)
    _write_backup(project, out_path, svg_files)
    if preview:
        try:
            from .preview_pptx import export_preview_pptx
            export_preview_pptx(project, stage=stage)
        except Exception:
            pass
    _preserve_notes(project, out_path)
    return out_path


def validate_pptx(path: Path | str) -> tuple[bool, list[str]]:
    deck = Path(path)
    errors: list[str] = []
    if not deck.exists():
        return False, [f"Missing PPTX: {deck}"]
    try:
        with zipfile.ZipFile(deck) as zf:
            names = set(zf.namelist())
            if "[Content_Types].xml" not in names:
                errors.append("Missing [Content_Types].xml")
            slide_names = [name for name in names if re.match(r"ppt/slides/slide\d+\.xml$", name)]
            if not slide_names:
                errors.append("No slide XML files found")
            native_shape_count = 0
            picture_count = 0
            for slide_name in slide_names:
                xml = zf.read(slide_name).decode("utf-8", errors="ignore")
                native_shape_count += xml.count("<p:sp>")
                picture_count += xml.count("<p:pic>")
            if native_shape_count == 0:
                errors.append("No native editable shape/text objects detected")
            if picture_count >= max(1, len(slide_names)) and native_shape_count == 0:
                errors.append("Deck appears to be image-only")
    except zipfile.BadZipFile:
        errors.append("File is not a valid ZIP/PPTX package")
    return not errors, errors


def pptx_text(path: Path | str) -> str:
    slides = extract_pptx_slide_text(path)
    return "\n\n".join(f"## Slide {idx}\n{text}" for idx, text in slides)


def pptx_notes(path: Path | str) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("Notes extraction requires python-pptx. Install with: python -m pip install -e .") from exc

    prs = Presentation(str(path))
    sections: list[str] = []
    for index, slide in enumerate(prs.slides, start=1):
        if not slide.has_notes_slide:
            continue
        text = slide.notes_slide.notes_text_frame.text.strip()
        if text:
            sections.append(f"## Slide {index}\n{text}")
    return "\n\n".join(sections)


def _embed_slide_notes(slide, note_text: str) -> None:
    note = note_text.strip()
    if not note:
        return
    _embed_rich_notes(slide, note)


def _embed_rich_notes(slide, text: str) -> None:
    """Embed notes with basic markdown formatting (bold, italic, bullets)."""
    from pptx.util import Pt

    try:
        from pptx.dml.color import RGBColor
    except ImportError:
        RGBColor = None

    notes_tf = slide.notes_slide.notes_text_frame

    if not _has_markdown(text):
        notes_tf.text = text
        return

    lines = text.split("\n")
    first = True
    for line in lines:
        if first:
            paragraph = notes_tf.paragraphs[0]
            first = False
        else:
            paragraph = notes_tf.add_paragraph()

        stripped = line.strip()
        if stripped.startswith("- ") or stripped.startswith("* "):
            paragraph.level = 1
            stripped = stripped[2:]

        for run_text, bold, italic in _parse_runs(stripped):
            run = paragraph.add_run()
            run.text = run_text
            if bold:
                run.font.bold = True
            if italic:
                run.font.italic = True


def _has_markdown(text: str) -> bool:
    return bool(re.search(r"\*\*(.+?)\*\*|\*(.+?)\*|^[-*]\s+", text, re.MULTILINE))


def _parse_runs(text: str) -> list[tuple[str, bool, bool]]:
    """Parse inline markdown into (text, bold, italic) run tuples."""
    runs: list[tuple[str, bool, bool]] = []
    pos = 0
    pattern = re.compile(r"(\*\*(.+?)\*\*|\*(.+?)\*)")

    for match in pattern.finditer(text):
        if match.start() > pos:
            runs.append((text[pos:match.start()], False, False))
        if match.group(2):
            runs.append((match.group(2), True, False))
        elif match.group(3):
            runs.append((match.group(3), False, True))
        pos = match.end()

    if pos < len(text):
        runs.append((text[pos:], False, False))

    if not runs:
        runs.append((text, False, False))

    return runs


def _read_project_notes(project: Path, slide_count: int) -> list[str]:
    notes = [""] * slide_count
    notes_dir = project / "notes"
    total = notes_dir / "total.md"
    if total.exists():
        notes = _parse_total_notes(total.read_text(encoding="utf-8"), slide_count)
    if notes_dir.exists():
        for note_file in sorted(notes_dir.glob("slide*.md")):
            match = re.search(r"slide[_-]?0*(\d+)", note_file.stem, re.IGNORECASE)
            if not match:
                continue
            index = int(match.group(1))
            if 1 <= index <= slide_count:
                notes[index - 1] = note_file.read_text(encoding="utf-8").strip()
    return notes


def _parse_total_notes(text: str, slide_count: int) -> list[str]:
    notes = [""] * slide_count
    current: int | None = None
    sections_found = False
    for line in text.splitlines():
        match = re.match(r"^\s*#{1,6}\s*Slide\s+0*(\d+)\b[:\-\s]*(.*)$", line, re.IGNORECASE)
        if match:
            index = int(match.group(1))
            current = index if 1 <= index <= slide_count else None
            sections_found = True
            heading_tail = match.group(2).strip()
            if current and heading_tail:
                notes[current - 1] += heading_tail + "\n"
            continue
        if current:
            notes[current - 1] += line + "\n"
    if sections_found:
        return [note.strip() for note in notes]

    chunks = [chunk.strip() for chunk in re.split(r"\n\s*---+\s*\n", text.strip()) if chunk.strip()]
    if len(chunks) > 1:
        for index, chunk in enumerate(chunks[:slide_count]):
            notes[index] = chunk
        return notes
    if slide_count == 1:
        notes[0] = text.strip()
    return notes


def _write_backup(project: Path, out_path: Path, svg_files: list[Path]) -> None:
    backup_dir = ensure_dir(project / "backup" / out_path.stem)
    if (project / "svg_output").exists():
        copy_tree(project / "svg_output", backup_dir / "svg_output")
    if (project / "svg_final").exists():
        copy_tree(project / "svg_final", backup_dir / "svg_final")
    manifest = backup_dir / "README.md"
    manifest.write_text(
        f"# Export Backup\n\nMain deck: `{out_path}`\n\nSVG files archived: {len(svg_files)}\n",
        encoding="utf-8",
    )


def _preserve_notes(project: Path, out_path: Path) -> None:
    notes = project / "notes" / "total.md"
    if notes.exists():
        shutil.copy2(notes, out_path.with_name(out_path.stem + "_notes.md"))
