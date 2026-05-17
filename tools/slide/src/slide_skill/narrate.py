"""TTS audio narration from speaker notes — supports edge-tts and MiMo TTS."""

from __future__ import annotations

import asyncio
from pathlib import Path

from .util import ensure_dir


async def _generate_audio_edge(text: str, output_path: Path, voice: str = "zh-CN-XiaoxiaoNeural") -> Path:
    import edge_tts

    communicate = edge_tts.Communicate(text, voice)
    await communicate.save(str(output_path))
    return output_path


def _generate_audio_mimo(
    text: str,
    output_path: Path,
    *,
    voice: str = "冰糖",
    style: str = "",
    voice_clone_sample: Path | None = None,
    voice_design_prompt: str | None = None,
) -> Path:
    from .mimo_tts import generate_audio

    if voice_clone_sample:
        model = "mimo-v2.5-tts-voiceclone"
    elif voice_design_prompt:
        model = "mimo-v2.5-tts-voicedesign"
    else:
        model = "mimo-v2.5-tts"

    return generate_audio(
        text,
        output_path,
        voice=voice,
        style=style,
        model=model,
        voice_clone_sample=voice_clone_sample,
        voice_design_prompt=voice_design_prompt,
    )


def narrate_slide(
    text: str,
    output_path: Path,
    *,
    voice: str = "zh-CN-XiaoxiaoNeural",
    engine: str = "edge-tts",
    style: str = "",
    voice_clone_sample: Path | None = None,
    voice_design_prompt: str | None = None,
) -> Path:
    """Generate a single audio file from text using the selected TTS engine."""
    output_path.parent.mkdir(parents=True, exist_ok=True)

    if engine == "mimo":
        return _generate_audio_mimo(
            text,
            output_path,
            voice=voice,
            style=style,
            voice_clone_sample=voice_clone_sample,
            voice_design_prompt=voice_design_prompt,
        )

    return asyncio.run(_generate_audio_edge(text, output_path, voice))


def narrate_project(
    project_path: Path | str,
    voice: str = "zh-CN-XiaoxiaoNeural",
    *,
    engine: str = "edge-tts",
    style: str = "",
    voice_clone_sample: Path | None = None,
    voice_design_prompt: str | None = None,
) -> list[Path]:
    """Generate per-slide audio from speaker notes and embed into PPTX.

    Returns list of generated audio file paths.
    """
    from .exporter import _read_project_notes
    from .project import load_project

    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("TTS requires python-pptx. Install with: python -m pip install -e .") from exc

    if engine == "edge-tts":
        try:
            import edge_tts
        except ImportError as exc:
            raise RuntimeError("edge-tts engine requires edge-tts. Install with: pip install -e .[audio]") from exc

    project = Path(project_path)
    meta = load_project(project)

    exports = sorted((project / "exports").glob("*.pptx"), key=lambda p: p.stat().st_mtime, reverse=True)
    if not exports:
        raise FileNotFoundError("No exported PPTX found. Run export first.")
    deck_path = exports[0]

    audio_dir = ensure_dir(project / "audio")

    prs = Presentation(str(deck_path))
    slide_count = len(prs.slides)
    notes = _read_project_notes(project, slide_count)

    ext = "wav" if engine == "mimo" else "mp3"
    generated: list[Path] = []

    for idx, note_text in enumerate(notes, start=1):
        text = note_text.strip()
        if not text:
            continue

        audio_path = audio_dir / f"slide_{idx:02d}.{ext}"
        narrate_slide(
            text,
            audio_path,
            voice=voice,
            engine=engine,
            style=style,
            voice_clone_sample=voice_clone_sample,
            voice_design_prompt=voice_design_prompt,
        )
        generated.append(audio_path)

        _embed_audio(prs.slides[idx - 1], audio_path)

    prs.save(deck_path)
    return generated


def _embed_audio(slide, audio_path: Path) -> None:
    """Embed an audio file into a PPTX slide."""
    from pptx.oxml.ns import qn
    from lxml import etree

    slide_part = slide.part
    audio_path_obj = Path(audio_path)

    with open(audio_path_obj, "rb") as f:
        audio_blob = f.read()

    suffix = audio_path_obj.suffix.lower()
    mime = "audio/wav" if suffix == ".wav" else "audio/mp3"

    rel = slide_part.relate_to(
        audio_blob,
        mime,
        is_external=False,
    )

    sp_tree = slide._element.find(qn("p:cSld")).find(qn("p:spTree"))

    pic = etree.SubElement(sp_tree, qn("p:pic"))
    nvPicPr = etree.SubElement(pic, qn("p:nvPicPr"))
    cNvPr = etree.SubElement(nvPicPr, qn("p:cNvPr"))
    cNvPr.set("id", "0")
    cNvPr.set("name", "Audio")
    nvPicPr_child = etree.SubElement(nvPicPr, qn("p:cNvPicPr"))
    nvPicPr_child.set("preferRelativeResize", "0")
    nvPr = etree.SubElement(nvPicPr, qn("p:nvPr"))

    a_audioFile = etree.SubElement(nvPr, qn("a:audioFile"))
    a_audioFile.set(qn("r:link"), rel.rId)

    blipFill = etree.SubElement(pic, qn("p:blipFill"))

    extent = etree.SubElement(pic, qn("p:extent"))
    cx = etree.SubElement(extent, qn("a:cx"))
    cx.text = "0"


def list_available_voices(locale: str = "", engine: str = "edge-tts") -> list[str]:
    """List available voices for the selected engine, optionally filtered by locale prefix."""
    if engine == "mimo":
        from .mimo_tts import list_mimo_voices
        voices = list_mimo_voices()
        result = [v["id"] for v in voices]
        if locale:
            result = [v for v in result if locale.lower() in v.lower()]
        return sorted(result)

    import edge_tts

    voices = asyncio.run(edge_tts.list_voices())
    result = []
    for v in voices:
        short = v["ShortName"]
        if not locale or short.startswith(locale):
            result.append(short)
    return sorted(result)
